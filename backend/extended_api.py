#!/usr/bin/env python3
from common.helpers import _aggregate_compute_results, _execute_common_step, _execute_compute_step, _execute_single_step, _execute_step, _finalize_common_operation, _finalize_results, _finalize_step_results, _initialize_compute_context, _prepare_common_context, _prepare_context, _prepare_step_context, _notify_progress


def _run_test_gate_simple(test_case: dict, config: dict) -> dict:
    """简化版测试门控检查。"""
    return {
        "passed": True,
        "test_case": test_case.get("name", ""),
        "score": test_case.get("score", 0)
    }

def _prepare_test_config(request_data: dict) -> dict:
    """简化版准备测试配置。"""
    return {
        "test_cases": request_data.get("test_cases", []),
        "threshold": request_data.get("threshold", 0.8)
    }



from typing import Any, Optional, Union, List, Dict, Tuple, Callable, Set, TypeVar, Generic, Iterator, Sequence, Mapping, Iterable, Awaitable, Coroutine, Type
from dataclasses import dataclass, field
from enum import Enum, auto
from datetime import datetime
import asyncio
from typing import Any, Optional, Union, List, Dict, Tuple, Callable, Set, TypeVar, Generic
from dataclasses import dataclass, field
from enum import Enum, auto
from datetime import datetime
"""Platform v9.0 Extended API - 研发增强/内容创作/运营分析/办公效率"""

import glob
import hashlib
import json
import logging
import os
import re
import socket
import subprocess
import threading
import time
import traceback
import uuid
from collections.abc import Callable
from datetime import datetime

from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel, Field

from common.auth import require_auth
from common.db import get_db, get_db_context
from common.llm import call_llm, call_llm_async, parse_llm_json, _safe_exc_msg
try:
    from prd_engine import stream_llm_response  # 完整版
except ImportError:
    # 独立版兜底：简单 SSE 流式
    import json as _json

    def stream_llm_response(system_prompt, user_prompt, max_tokens, usage_key):
        from common.llm import call_llm_async
        from fastapi.responses import StreamingResponse

        async def _gen():
            try:
                full = await call_llm_async(system_prompt, user_prompt, max_tokens=max_tokens)
                yield "data: " + _json.dumps({"delta": full}) + "\n\n"
                yield "data: " + _json.dumps({"done": True, "full": full}) + "\n\n"
            except Exception as e:
                yield "data: " + _json.dumps({"error": str(e)}) + "\n\n"

        return StreamingResponse(_gen(), media_type="text/event-stream")

from task_queue import create_task, register_handler

logger = logging.getLogger(__name__)

router = APIRouter()


# ══════════════════════════════════════════════════════════════
# Phase 2: 研发增强
# ══════════════════════════════════════════════════════════════


class CodeGenRequest(BaseModel):
    language: str = "python"
    prompt: str
    model: str = ""


class CodeReviewRequest(BaseModel):
    language: str = "python"
    code: str
    model: str = ""
    stream: bool = False


class CodeImproveRequest(BaseModel):
    """根据代码审查意见修改代码"""

    language: str = "python"
    code: str
    review: str
    model: str = ""


class PipelineCreate(BaseModel):
    name: str
    description: str = ""
    type: str = "ci"
    config: dict = {}


class DeployRequest(BaseModel):
    """一键部署请求：代码落盘 + 构建镜像 + 沙箱容器运行"""

    name: str
    language: str = "python"
    code: str
    requirement_id: str = ""


# ── 部署流水线（真实执行：podman 构建镜像 → 启动沙箱容器） ──────
ARTIFACTS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "artifacts")
DEPLOY_BASE_PORT = 18080
_PODMAN_LOCK = threading.Lock()

# 代码 import 顶层模块 → pip 包名（用于生成 requirements.txt）
_PIP_PACKAGES = {
    "fastapi": "fastapi",
    "uvicorn": "uvicorn",
    "flask": "flask",
    "requests": "requests",
    "httpx": "httpx",
    "pydantic": "pydantic",
    "sqlalchemy": "sqlalchemy",
    "numpy": "numpy",
    "pandas": "pandas",
    "aiohttp": "aiohttp",
    "redis": "redis",
    "openai": "openai",
    "dotenv": "python-dotenv",
    "bs4": "beautifulsoup4",
    "PIL": "pillow",
    "yaml": "pyyaml",
    "flask_cors": "flask-cors",
    "jwt": "pyjwt",
    "celery": "celery",
    "django": "django",
    "click": "click",
}


def _detect_python_deps(code: str) -> list:
    """从代码 import 语句提取 pip 依赖（Web 服务基础依赖兜底）。"""
    deps = {"fastapi", "uvicorn"}
    for m in re.finditer(r"^\s*(?:from|import)\s+([a-zA-Z_][\w\.]*)", code, re.M):
        top = m.group(1).split(".")[0]
        if top in _PIP_PACKAGES:
            deps.add(_PIP_PACKAGES[top])
    return sorted(deps)


def _detect_project_type(project_dir: str) -> dict:
    """检测项目技术栈，返回通用流水线信息。

    {lang: python|node|go|docker, entry: 入口文件, test_file: 测试文件名,
     test_cmd: 测试命令列表, container_port: 容器内服务端口, has_dockerfile: 是否已有 Dockerfile}
    - 语言优先按工程清单文件推断（package.json → node，go.mod → go，否则 python）
    - 已有 Dockerfile 的项目直接复用（尊重用户容器化配置，最通用），仅读 EXPOSE 端口
    """
    files = set()
    for root, dirs, fs in os.walk(project_dir):
        dirs[:] = [
            d
            for d in dirs
            if d
            not in (
                "node_modules",
                ".git",
                "__pycache__",
                ".venv",
                "venv",
                "dist",
                "build",
                ".next",
                "coverage",
                "target",
            )
        ]
        for f in fs:
            files.add(os.path.relpath(os.path.join(root, f), project_dir))
    info = {
        "lang": "python",
        "entry": "main.py",
        "test_file": "test_main.py",
        "test_cmd": ["pytest", "-q", "--tb=short", "test_main.py"],
        "container_port": 8000,
        "has_dockerfile": False,
    }
    if "package.json" in files:
        info["lang"] = "node"
        info["entry"] = next(
            (e for e in ("server.js", "index.js", "app.js", "main.js", "src/index.js") if e in files), "server.js"
        )
        info["test_file"] = "test_api.test.js"
        info["test_cmd"] = ["npm", "test"]
    elif "go.mod" in files:
        info["lang"] = "go"
        info["entry"] = "main.go"
        info["test_file"] = "main_test.go"
        info["test_cmd"] = ["go", "test", "./..."]
    else:
        # python 工程：入口缺失 main.py 时按典型入口名/首个非测试 .py 推断
        if "main.py" not in files:
            pys = sorted(
                f for f in files if f.endswith(".py") and not f.startswith("test_") and not f.endswith("_test.py")
            )
            info["entry"] = next(
                (c for c in ("app.py", "server.py", "api.py", "index.py") if c in pys), pys[0] if pys else "main.py"
            )
    # 已有 Dockerfile 的项目直接复用（尊重用户容器化配置，最通用）
    if "Dockerfile" in files or "dockerfile" in files:
        info["has_dockerfile"] = True
        try:
            with open(os.path.join(project_dir, "Dockerfile"), encoding="utf-8") as f:
                for line in f:
                    m = re.search(r"EXPOSE\s+(\d+)", line)
                    if m:
                        info["container_port"] = int(m.group(1))
                        break
        except Exception:
            pass
    return info


def _gen_dockerfile(lang: str = "python", include_tests: bool = False, entry: str = "main.py") -> str:
    """按语言生成服务 Dockerfile（容器内固定 8000 端口）。

    include_tests=True 时生成测试镜像 Dockerfile：额外复制测试文件并安装测试依赖，
    用于自动化测试门禁在容器内真实执行用例。已有 Dockerfile 的项目由调用方直接复用。
    """
    if lang == "node":
        base = (
            "FROM node:20-alpine\n"
            "WORKDIR /app\n"
            "COPY package*.json .\n"
            "RUN npm install --registry=https://registry.npmmirror.com || npm install\n"
            "COPY . .\n"
            "ENV PORT=8000\n"
            "EXPOSE 8000\n"
        )
        if include_tests:
            return base + 'CMD ["npm", "test"]\n'
        return base + f'CMD ["node", "{entry}"]\n'
    if lang == "go":
        if include_tests:
            return (
                "FROM golang:1.22-alpine\n"
                "WORKDIR /app\n"
                "COPY go.mod go.sum* .\n"
                "RUN go mod download || true\n"
                "COPY . .\n"
                'CMD ["go", "test", "./..."]\n'
            )
        return (
            "FROM golang:1.22-alpine AS build\n"
            "WORKDIR /app\n"
            "COPY go.mod go.sum* .\n"
            "RUN go mod download || true\n"
            "COPY . .\n"
            "RUN CGO_ENABLED=0 go build -o /app/server .\n"
            "FROM alpine:3.19\n"
            "WORKDIR /app\n"
            "COPY --from=build /app/server /app/server\n"
            "ENV PORT=8000\n"
            "EXPOSE 8000\n"
            'CMD ["/app/server"]\n'
        )
    # python 默认（COPY . . 全量复制，支持多文件工程；入口文件按 entry 指定）
    if include_tests:
        return (
            "FROM python:3.11-slim\n"
            "WORKDIR /app\n"
            "COPY . .\n"
            "RUN pip install --no-cache-dir -r requirements.txt pytest\n"
            "EXPOSE 8000\n"
            f'CMD ["python", "{entry}"]\n'
        )
    return (
        "FROM python:3.11-slim\n"
        "WORKDIR /app\n"
        "COPY . .\n"
        "RUN pip install --no-cache-dir -r requirements.txt\n"
        "EXPOSE 8000\n"
        f'CMD ["python", "{entry}"]\n'
    )


def _find_free_port() -> int:
    """从 DEPLOY_BASE_PORT 起扫描空闲端口（部署用）。"""
    with _PODMAN_LOCK:
        for port in range(DEPLOY_BASE_PORT, DEPLOY_BASE_PORT + 300):
            with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
                s.settimeout(0.3)
                if s.connect_ex(("127.0.0.1", port)) != 0:
                    return port
    return DEPLOY_BASE_PORT


def _update_run_log(run_id: str, log: str) -> None:
    """更新流水线运行日志（渐进式写入）。"""
    conn = get_db()
    try:
        conn.execute("UPDATE pipeline_runs SET log=? WHERE id=?", (log, run_id))
        conn.commit()
    finally:
        conn.close()


def _finish_run(run_id: str, pid: str, status: str, log: str) -> None:
    """结束运行：写状态/日志并同步流水线状态。"""
    finished = datetime.now().isoformat()
    conn = get_db()
    try:
        conn.execute(
            "UPDATE pipeline_runs SET status=?, log=?, finished_at=? WHERE id=?",
            (status, log, finished, run_id),
        )
        conn.execute(
            "UPDATE pipelines SET status=?, last_run=? WHERE id=?",
            (status, finished, pid),
        )
        conn.commit()
    finally:
        conn.close()


def _safe_slug(name: str) -> str:
    """容器/镜像名安全化：Docker 标签仅允许 [a-zA-Z0-9][a-zA-Z0-9_.-]*，
    中文/特殊字符名用 md5 摘要后缀生成合法 slug（原名保留用于展示）。"""
    if re.fullmatch(r"[a-zA-Z0-9][a-zA-Z0-9_.-]*", name):
        return name
    digest = hashlib.md5(name.encode("utf-8")).hexdigest()[:8]
    return f"svc-{digest}"


_INFRA_ERROR_MARKERS = (
    "cannot connect to podman",
    "connection refused",
    "no space left",
    "command not found",
    "permission denied",
    "cannot connect to the docker daemon",
    "docker daemon is not running",
    "executable file not found",
)


def _is_infra_error(msg: str) -> bool:
    """是否为基础设施故障（Docker/Podman 环境问题）——此类错误 AI 改代码无法解决，跳过修复。"""
    low = (msg or "").lower()
    return any(m in low for m in _INFRA_ERROR_MARKERS)



def _ensure_dep_redis(cfg: dict, container_name: str, net: str, append, step_run) -> tuple:
    """确保 Redis 依赖容器运行（幂等复用）。返回 (env_flags, ok, err)。"""
    env_flags: list = []
    dep = f"{container_name}-redis"
    ok, out = step_run(["podman", "inspect", "--format", "{{.State.Running}}", dep], timeout=30)
    if ok and out.strip() == "true":
        append(f"  - 依赖: Redis 容器 {dep} 已运行，复用 ✓")
    else:
        step_run(["podman", "rm", "-f", dep], timeout=30)
        append(f"  - 依赖: 启动 Redis 容器 {dep} …")
        ok, out = step_run(
            ["podman", "run", "-d", "--name", dep, "--network", net, "--network-alias", "redis",
             "docker.io/library/redis:7-alpine"],
            timeout=300,
        )
        if not ok:
            return [], False, f"Redis 依赖容器启动失败: {out[-300:]}"
    env_flags += ["-e", "REDIS_URL=redis://redis:6379/0"]
    return env_flags, True, ""


def _ensure_dep_mysql(cfg: dict, container_name: str, net: str, append, step_run) -> tuple:
    """确保 MySQL 依赖容器运行并就绪（幂等复用）。返回 (env_flags, ok, err)。"""
    env_flags: list = []
    dep = f"{container_name}-mysql"
    pw = cfg.get("mysql_password", "platform123")
    db = cfg.get("mysql_database", "platform")
    img = cfg.get("mysql_image", "docker.io/library/mysql:8")
    ok, out = step_run(["podman", "inspect", "--format", "{{.State.Running}}", dep], timeout=30)
    if ok and out.strip() == "true":
        ok2, _ = step_run(["podman", "exec", dep, "mysqladmin", "ping", "-uroot", f"-p{pw}", "--silent"], timeout=10)
        if ok2:
            append(f"  - 依赖: MySQL 容器 {dep} 已运行且就绪，复用 ✓")
        else:
            append(f"  - 依赖: MySQL 容器 {dep} 未就绪，重建 …")
            ok = False
    else:
        ok = False
    if not ok:
        step_run(["podman", "rm", "-f", dep], timeout=30)
        append(f"  - 依赖: 启动 MySQL 容器 {dep} …（首次拉取镜像较慢）")
        ok, out = step_run(
            ["podman", "run", "-d", "--name", dep, "--network", net, "--network-alias", "mysql",
             "-e", f"MYSQL_ROOT_PASSWORD={pw}", "-e", f"MYSQL_DATABASE={db}", img],
            timeout=600,
        )
        if not ok:
            return [], False, f"MySQL 依赖容器启动失败: {out[-300:]}"
        append("  - 依赖: 等待 MySQL 初始化（约 20-40s）…")
        ready = False
        for _ in range(60):
            time.sleep(2)
            ok2, _ = step_run(["podman", "exec", dep, "mysqladmin", "ping", "-uroot", f"-p{pw}", "--silent"], timeout=10)
            if ok2:
                ready = True
                break
        if not ready:
            return [], False, "MySQL 依赖容器初始化超时"
    env_flags += ["-e", f"DATABASE_URL=mysql+aiomysql://root:{pw}@mysql:3306/{db}?charset=utf8mb4"]
    return env_flags, True, ""

def _prepare_dependencies(cfg, container_name, append, step_run) -> tuple:  # noqa: C901
    """准备自定义网络 + 依赖容器（Redis/MySQL），已存在且健康则复用（幂等）。

    podman 5.x 已移除 --link，改用自定义网络 + --network-alias（同网络容器名/别名可直接解析）。
    返回 (net, env_flags, ok, err)。测试门禁与部署共用，避免重复重建依赖。
    """
    deps = cfg.get("dependencies") or {}
    net = f"{container_name}-net"
    env_flags: list = []
    # 无论是否有依赖，都确保自定义网络存在（幂等）：测试门禁与部署统一使用该网络，
    # 避免无依赖场景下 podman run --network xxx-net 报 network not found
    ok, out = step_run(["podman", "network", "exists", net], timeout=30)
    if not ok:
        ok, out = step_run(["podman", "network", "create", net], timeout=30)
        if not ok and "already exists" not in (out or "").lower():
            return "", [], False, f"创建网络失败: {out[-300:]}"
    if deps:
        if deps.get("redis"):
            flags, ok, err = _ensure_dep_redis(deps, container_name, net, append, step_run)
            if not ok:
                return "", [], False, err
            env_flags += flags
        if deps.get("mysql"):
            flags, ok, err = _ensure_dep_mysql(deps, container_name, net, append, step_run)
            if not ok:
                return "", [], False, err
            env_flags += flags
    return net, env_flags, True, ""


def _deploy_once(name, project_dir, port, image_tag, container_name, append, step_run, cfg=None) -> tuple:
    """单轮部署：构建镜像 → 启动依赖容器 → 启动服务容器 → 健康检查。返回 (ok, info)。"""
    cfg = cfg or {}
    # 阶段 1.5：外部依赖容器（Redis/MySQL）。幂等准备：已运行则复用，避免测试门禁后重复初始化
    deps = cfg.get("dependencies") or {}
    net, env_flags, ok, err = _prepare_dependencies(cfg, container_name, append, step_run)
    if not ok:
        return False, err

    # 阶段 2：构建镜像
    append(f"  - 构建镜像: podman build -t {image_tag} …（首次拉取基础镜像较慢）")
    ok, out = step_run(["podman", "build", "-t", image_tag, project_dir])
    if out:
        lines = out.splitlines()
        tail = lines[-3:] if len(lines) > 3 else lines
        append("      " + "\n      ".join(tail))
    if not ok:
        return False, f"镜像构建失败: {out[-400:]}"
    append("  - 构建镜像: 完成 ✓")
    # 阶段 3：启动沙箱容器（先清理同名旧容器，支持修复重部署）
    step_run(["podman", "rm", "-f", container_name], timeout=30)
    container_port = cfg.get("container_port") or _detect_project_type(project_dir)["container_port"]
    append(
        f"  - 启动容器: podman run -d --name {container_name} -p {port}:{container_port}"
        + (f" --network {net}" if deps else "")
    )
    cmd = ["podman", "run", "-d", "--name", container_name, "-p", f"{port}:{container_port}"]
    if deps:
        cmd += ["--network", net]
    cmd += env_flags
    cmd += [image_tag]
    ok, out = step_run(cmd)
    if out:
        append("      " + out)
    if not ok:
        return False, f"容器启动失败: {out[-400:]}"
    # 阶段 4：健康检查（轮询端口，最多 60s）
    append(f"  - 健康检查: http://localhost:{port}/ …")
    code = "000"
    for _ in range(30):
        time.sleep(2)
        try:
            r = subprocess.run(
                ["curl", "-s", "-o", "/dev/null", "-w", "%{http_code}", f"http://localhost:{port}/"],
                capture_output=True,
                text=True,
                timeout=5,
            )
            code = r.stdout.strip()
            if code and code != "000":
                return True, code
        except Exception:
            pass
    return False, "健康检查未通过，服务可能启动失败（查看容器日志定位问题）"


def _register_sandbox(slug, port, project_dir, image_tag, cfg, display_name=None) -> None:
    """把部署服务注册到沙箱管理（可在沙箱页停止/删除/查看日志）。slug 为容器安全名，display_name 为展示名。"""
    name = display_name or slug
    now = datetime.now().isoformat()
    conn = get_db()
    try:
        conn.execute(
            "INSERT INTO sandbox_projects (id, name, status, port, project_dir, image, ports, config, created_at, updated_at) "
            "VALUES (?,?,?,?,?,?,?,?,?,?) ON CONFLICT(id) DO UPDATE SET status=?, port=?, project_dir=?, image=?, ports=?, config=?, updated_at=?",
            (
                f"deploy-{slug}",
                name,
                "running",
                port,
                project_dir,
                image_tag,
                json.dumps([port]),
                json.dumps(cfg, ensure_ascii=False),
                now,
                now,
                "running",
                port,
                project_dir,
                image_tag,
                json.dumps([port]),
                json.dumps(cfg, ensure_ascii=False),
                now,
            ),
        )
        conn.commit()
    finally:
        conn.close()


def _extract_code_block(text: str) -> str:
    """从 LLM 输出提取完整代码（清洗 markdown 围栏）。

    - 优先提取所有 ``` 围栏代码块，取最长的一个（容忍前置解释文字/前导换行）
    - 退化场景：只有开头围栏无闭合围栏（LLM 偶发截断），剥离首个围栏行后返回
    - 无围栏时若全文无明显解释性文字则原样返回
    """
    text = (text or "").strip()
    if not text:
        return ""
    blocks = re.findall(r"```[a-zA-Z]*\s*\n(.*?)```", text, re.DOTALL)
    if blocks:
        return max(blocks, key=len).strip()
    # 退化：只有开头围栏、无闭合围栏（LLM 输出被截断时常见）→ 剥离首个围栏行
    cleaned = re.sub(r"^```[a-zA-Z]*\s*\n?", "", text).strip()
    if cleaned != text:
        # 剥离后仍带中文解释文字且无代码特征 → 视为无效输出
        if (
            re.search(r"[，。；：、]|以下是|修复建议|问题|错误|请提供|您好|需要", cleaned[:200])
            and "def " not in cleaned[:200]
            and "import " not in cleaned[:200]
            and "class " not in cleaned[:200]
        ):
            return ""
        return cleaned
    # 无明显解释性文字才视为纯代码（避免把 LLM 的说明文字写入 main.py）
    if re.search(r"[，。；：、]|以下是|修复建议|问题|错误|请提供|您好|需要", text[:200]):
        return ""
    return text


def _fix_rounds(pid, run_id, cfg, log, append, step_run, initial_error, max_rounds=3) -> bool:
    """AI 诊断修复循环：错误日志 → LLM 修改入口文件 → 重建 → 重启 → 健康检查。

    按语言（python/node/go/docker）生成修复提示词，每轮把最新错误日志喂给 LLM，
    最多 max_rounds 轮；成功返回 True。
    """
    name = cfg["service_name"]
    slug = _safe_slug(name)
    project_dir = cfg["project_dir"]
    port = cfg["port"]
    image_tag = f"app-{slug}"
    container_name = f"sandbox-{slug}"
    lang = cfg.get("language") or _detect_project_type(project_dir)["lang"]
    entry = cfg.get("entry") or "main.py"
    main_file = os.path.join(project_dir, entry)
    last_error = initial_error
    lang_prompts = {
        "python": (
            "你是一个资深的 Python 开发工程师与 SRE。根据下面的部署错误日志和当前入口文件代码，"
            "定位问题根因并输出修复后的完整入口文件（必须是可直接运行的 Web 服务，监听 0.0.0.0:8000，"
            "根路径 / 必须返回 200）。只返回修复后的完整代码，放在 ```python 代码块中，不要任何解释文字。"
        ),
        "node": (
            "你是一个资深的 Node.js 开发工程师与 SRE。根据下面的部署错误日志和当前入口文件代码，"
            "定位问题根因并输出修复后的完整入口文件（必须是可直接运行的 HTTP 服务，监听 0.0.0.0:8000，"
            "根路径 / 必须返回 200）。只返回修复后的完整代码，放在 ```javascript 代码块中，不要任何解释文字。"
        ),
        "go": (
            "你是一个资深的 Go 开发工程师与 SRE。根据下面的部署错误日志和当前入口文件代码，"
            "定位问题根因并输出修复后的完整入口文件（必须是可直接运行的 HTTP 服务，监听 0.0.0.0:8000，"
            "根路径 / 必须返回 200）。只返回修复后的完整代码，放在 ```go 代码块中，不要任何解释文字。"
        ),
    }
    sys_prompt = lang_prompts.get(lang, lang_prompts["python"])
    for round_no in range(1, max_rounds + 1):
        append(f"── 第 {round_no}/{max_rounds} 轮 AI 诊断修复 ──")
        # 收集诊断信息：上轮错误 + 容器日志
        _, clogs = step_run(["podman", "logs", "--tail", "80", container_name], timeout=10)
        diag = last_error
        if clogs:
            diag += "\n【容器日志】\n" + clogs[-2500:]
        try:
            with open(main_file, encoding="utf-8") as f:
                code = f.read()
        except Exception as e:
            append(f"  - ❌ 读取代码失败: {e}")
            return False
        append("  - LLM 诊断失败原因并生成修复代码（约 10-60 秒）…")
        prompt = f"【部署错误日志】\n{diag}\n\n【当前 {entry}】\n{code}"
        try:
            fix = call_llm(sys_prompt, prompt)
        except Exception as e:
            append(f"  - ❌ LLM 调用失败: {e}（可在系统配置-模型列表中设置模型 API Key）")
            return False
        fixed = _extract_code_block(fix)
        if not fixed:
            # LLM 可能返回了解释文字/空内容：强制要求只输出代码块，再试一次
            append("  - ⚠ LLM 未按格式输出代码，要求重新输出…")
            try:
                fix2 = call_llm(
                    sys_prompt,
                    prompt + "\n\n（你上一次没有输出代码块。请只输出代码块围栏包裹的完整代码，不要任何解释文字。）",
                )
            except Exception as e:
                append(f"  - ❌ LLM 重试调用失败: {e}")
                return False
            fixed = _extract_code_block(fix2)
            if not fixed:
                append("  - ❌ 未解析到修复代码，停止修复")
                return False
            append("  - 重新输出成功，继续修复…")
        with open(main_file, "w", encoding="utf-8") as f:
            f.write(fixed)
        append(f"  - 修复代码已落盘（{len(fixed)} 字节），重新构建部署…")
        ok, info = _deploy_once(name, project_dir, port, image_tag, container_name, append, step_run, cfg)
        if ok:
            append(f"  - 第 {round_no} 轮修复成功 ✓（HTTP {info}）访问地址: http://localhost:{port}")
            _register_sandbox(slug, port, project_dir, image_tag, cfg, display_name=name)
            return True
        if _is_infra_error(info):
            append("  - ⚠ 检测到基础设施故障（Docker/Podman 环境问题），AI 修复无法解决，停止修复")
            return False
        append(f"  - 第 {round_no} 轮仍失败: {info[-300:]}")
        last_error = info
    append("  - ❌ AI 修复达到轮次上限，请人工查看日志处理")
    return False


TEST_FIX_SYSTEM = (
    "你是一个资深的 Python 开发工程师。根据下面的构建/测试失败输出和当前 main.py 代码，"
    "定位问题根因并输出修复后的完整 main.py（必须是可直接运行的 Web 服务，监听 0.0.0.0:8000，"
    "提供 FastAPI 应用，可包含 Flask 等，但根路径 / 必须返回 200）。"
    "保持原有功能与接口不变，只修复导致失败的问题。"
    "只返回修复后的完整代码，放在 ```python 代码块中，不要任何解释文字。"
)

TEST_FILE_FIX_SYSTEM = (
    "你是一个资深的测试工程师。根据下面的测试执行/语法错误输出和当前测试文件内容，"
    "定位问题根因并输出修复后的完整 pytest 测试文件。"
    "保持测试意图与覆盖范围不变，只修复导致失败的问题（语法/import/断言等）。"
    "只返回修复后的完整代码，放在 ```python 代码块中，不要任何解释文字。"
)

PATCH_FIX_SYSTEM = (
    "你是一个资深的 Python 开发工程师。根据下面的失败输出，定位问题根因，"
    "输出修复代码所需的 unified diff（diff -u 格式）：\n"
    "1. 只输出 diff 本身，放在 ```diff 围栏中，不要输出完整文件，不要任何解释文字\n"
    "2. 文件头格式：--- a/文件名 和 +++ b/文件名（patch -p1 应用）\n"
    "3. 只包含修复所需的最小改动，保留原有缩进，diff 必须格式合法\n"
    "4. 若失败源于测试期望与实现不符（如外部依赖不可用），给出合理修复（mock 外部调用或调整逻辑），"
    "   不得删除接口或改变接口语义"
)

FUNCTION_FIX_SYSTEM = (
    "你是一个资深的 Python 开发工程师。根据下面的失败输出和当前函数代码，定位问题根因，"
    "只输出修复后的完整函数定义：\n"
    "1. 输出顺序：先按需输出修复后的路由装饰器行（@app.xxx(...)，仅当必须修改时，如 response_model "
    "   与返回结构冲突），紧接着输出 def 函数定义本身（从 def 行到函数结束），不要 import、其他函数或解释文字\n"
    "2. 保持既有参数不变，可新增可选参数（如 type: str = Query(None)）以满足测试用例的参数需求；"
    "   保留 4 空格缩进风格\n"
    "3. 若失败源于外部依赖不可用（如第三方 API Key 缺失/网络失败/下游 500），在函数内做合理降级"
    "   （返回友好错误码或 mock 数据），不要抛出未处理异常；注意路由装饰器声明的 response_model："
    "   返回的 dict 必须包含该模型所有必填字段（缺字段会触发 FastAPI ResponseValidationError 导致 500）\n"
    "4. 若提供了失败测试用例，降级返回的数据结构与状态码必须满足其断言"
    "   （如 data 下必须含 live/forecast 字段，或特定状态码），可用 mock/示例数据填充；"
    "   可通过异常信息区分场景（如 'not found' 返回 404、网络/Key 问题降级 200）；"
    "   若测试要求 data 为 list（如 type=forecast 返回 5 条列表）或 data 直接为 dict（如 type=live），"
    "   与 response_model 的 Dict 约束冲突时，必须修改装饰器（移除 response_model 或改为兼容模型）\n"
    "5. 修复后的函数必须完整不截断"
)


_TEST_PROMPTS = {
    "python": (
        "你是一个资深的测试工程师。根据需求测试用例和 main.py 代码，生成 pytest 测试文件 test_main.py：\n"
        "1. 使用真实 HTTP 测试：fixture 中 subprocess 启动 uvicorn main:app（随机端口如 8911），用 httpx 同步请求；\n"
        "   禁止使用 fastapi TestClient——若 main.py 接口为同步函数内部 asyncio.run()（AI 生成常见），TestClient 会报 Task attached to different loop\n"
        "2. 覆盖核心接口的正常路径、参数校验与错误分支，断言要合理；FastAPI 参数校验失败返回 422（断言以 422 为准，不要预期 400）\n"
        "3. 不依赖真实外部第三方 API（外部网络调用在服务端做降级或跳过）\n"
        "4. 禁止使用 unittest.mock.patch 对 main 模块做 mock（测试运行在独立子进程 uvicorn 中，patch 不生效）\n"
        "5. 全部使用同步调用，禁止 async/await 异步测试（容器未安装 pytest-asyncio）\n"
        "6. 测试文件必须自包含（只 import 标准库/httpx/pytest），通过环境变量连接依赖（REDIS_URL/DATABASE_URL 已注入）\n"
        "7. 只输出 ```python 围栏包裹的完整测试代码，不要任何解释文字，代码必须完整不截断\n"
        "8. 每个测试函数必须写 docstring，首行标注需求用例编号与标题（如 `TC-API-001: 搜索北京`），编号严格取自【需求测试用例】文档；需求用例缺失编号时写中文场景名（如 `搜索北京`）\n"
        "\n推荐测试骨架（可直接使用）：\n"
        "import os, subprocess, time, httpx, pytest\n"
        '@pytest.fixture(scope="module")\n'
        "def base_url():\n"
        '    port = os.environ.get("TEST_PORT", "8911")\n'
        '    proc = subprocess.Popen(["uvicorn", "main:app", "--host", "127.0.0.1", "--port", port],\n'
        "        env={**os.environ}, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)\n"
        "    for _ in range(40):\n"
        "        try:\n"
        '            httpx.get(f"http://127.0.0.1:{port}/docs", timeout=1)\n'
        "            break\n"
        "        except Exception:\n"
        "            time.sleep(0.5)\n"
        '    yield f"http://127.0.0.1:{port}"\n'
        "    proc.terminate()\n"
        '测试用例通过 httpx.get(base_url + "/路径", params=...) 调用并断言。'
    ),
    "node": (
        "你是一个资深的 Node.js 测试工程师。根据需求测试用例和入口文件代码，生成 node --test 测试文件 test_api.test.js：\n"
        "1. 使用真实 HTTP 测试：before 钩子中用 child_process.spawn 启动服务进程（node <entry>，环境变量 PORT 设为随机端口如 8911），"
        "用 Node 18+ 全局 fetch 同步请求\n"
        "2. 覆盖核心接口的正常路径、参数校验与错误分支，断言要合理\n"
        "3. 不依赖真实外部第三方 API（外部网络调用在服务端做降级或跳过）\n"
        "4. 测试文件必须自包含（只 require node:test/node:assert/child_process/process），"
        "通过环境变量连接依赖（REDIS_URL/DATABASE_URL 已注入）\n"
        "5. 只输出 ```javascript 围栏包裹的完整测试代码，不要任何解释文字，代码必须完整不截断\n"
        "\n推荐测试骨架（可直接使用）：\n"
        "const { test, before, after } = require('node:test');\n"
        "const assert = require('node:assert');\n"
        "const { spawn } = require('child_process');\n"
        "let proc, base;\n"
        "before(() => {\n"
        "  const port = process.env.TEST_PORT || '8911';\n"
        "  base = `http://127.0.0.1:${port}`;\n"
        "  proc = spawn(process.execPath, ['<entry>'], { env: { ...process.env, PORT: port }, stdio: 'ignore' });\n"
        "});\n"
        "after(() => proc && proc.kill());\n"
        "测试用例通过 fetch(base + '/路径') 调用并断言（注意等待服务就绪）。"
    ),
    "go": (
        "你是一个资深的 Go 测试工程师。根据需求测试用例和入口文件代码，生成 Go 测试文件 main_test.go：\n"
        "1. 使用真实 HTTP 测试：TestMain 中用 os/exec 启动服务进程（go run .，环境变量 PORT 设为随机端口如 8911），"
        "用 net/http 客户端同步请求\n"
        "2. 覆盖核心接口的正常路径、参数校验与错误分支，断言要合理\n"
        "3. 不依赖真实外部第三方 API（外部网络调用在服务端做降级或跳过）\n"
        "4. 测试函数命名 TestXxx(t *testing.T)，测试文件必须自包含（只 import 标准库），"
        "通过环境变量连接依赖（REDIS_URL/DATABASE_URL 已注入）\n"
        "5. 只输出 ```go 围栏包裹的完整测试代码，不要任何解释文字，代码必须完整不截断\n"
        "\n推荐测试骨架（可直接使用）：\n"
        "package main\n\n"
        'import (\n\t"net/http"\n\t"os"\n\t"os/exec"\n\t"testing"\n\t"time"\n)\n\n'
        'func TestMain(m *testing.M) {\n\tport := os.Getenv("TEST_PORT")\n\tif port == "" {\n\t\tport = "8911"\n\t}\n'
        '\tcmd := exec.Command("go", "run", ".")\n\tcmd.Env = append(os.Environ(), "PORT="+port)\n'
        "\t_ = cmd.Start()\n\ttime.Sleep(3 * time.Second)\n"
        "\tcode := m.Run()\n\t_ = cmd.Process.Kill()\n\tos.Exit(code)\n}\n\n"
        '测试用例通过 http.Get("http://127.0.0.1:" + port + "/路径") 调用并断言。'
    ),
}


def _validate_test_file(lang: str, content: str) -> tuple:
    """按语言校验测试代码语法，返回 (ok, err)。python 用 ast；node 用 node --check；go 用 gofmt -e。
    本地缺少对应运行时（node/go）时跳过校验（容器内执行时由测试命令兜底）。
    """
    import tempfile

    try:
        if lang == "python":
            import ast

            ast.parse(content)
            return True, ""
        suffix = {"node": ".js", "go": ".go"}.get(lang)
        if not suffix:
            return True, ""
        with tempfile.NamedTemporaryFile("w", suffix=suffix, delete=False, encoding="utf-8") as tf:
            tf.write(content)
            tmp = tf.name
        try:
            if lang == "node":
                r = subprocess.run(["node", "--check", tmp], capture_output=True, text=True, timeout=30)
                return (r.returncode == 0), (r.stderr or r.stdout or "").strip()[-500:]
            r = subprocess.run(["gofmt", "-e", tmp], capture_output=True, text=True, timeout=30)
            return (r.returncode == 0), (r.stderr or "").strip()[-500:]
        finally:
            os.unlink(tmp)
    except FileNotFoundError:
        return True, ""  # 本地无对应运行时：跳过语法校验
    except SyntaxError as e:
        return False, f"L{e.lineno}: {e.msg}"


def _fix_system(lang: str, kind: str) -> str:
    """按语言与修复目标返回 LLM 修复系统提示词（main=入口文件 / test_file=测试文件 / patch=最小补丁）。"""
    if lang == "node":
        if kind == "patch":
            return (
                "你是一个资深的 Node.js 开发工程师。根据下面的失败输出，定位问题根因，"
                "输出修复代码所需的 unified diff（diff -u 格式）：\n"
                "1. 只输出 diff 本身，放在 ```diff 围栏中，不要输出完整文件，不要任何解释文字\n"
                "2. 文件头格式：--- a/文件名 和 +++ b/文件名（patch -p1 应用）\n"
                "3. 只包含修复所需的最小改动，保留原有缩进，diff 必须格式合法\n"
                "4. 若失败源于测试期望与实现不符（如外部依赖不可用），给出合理修复（mock 外部调用或调整逻辑），"
                "不得删除接口或改变接口语义"
            )
        if kind == "test_file":
            return (
                "你是一个资深的测试工程师。根据下面的测试执行/语法错误输出和当前测试文件内容，"
                "定位问题根因并输出修复后的完整 node --test 测试文件。"
                "保持测试意图与覆盖范围不变，只修复导致失败的问题（语法/import/断言等）。"
                "只返回修复后的完整代码，放在 ```javascript 代码块中，不要任何解释文字。"
            )
        return (
            "你是一个资深的 Node.js 开发工程师。根据下面的构建/测试失败输出和当前入口文件代码，"
            "定位问题根因并输出修复后的完整入口文件（必须是可直接运行的 HTTP 服务，监听 0.0.0.0:8000，"
            "根路径 / 必须返回 200）。保持原有功能与接口不变，只修复导致失败的问题。"
            "只返回修复后的完整代码，放在 ```javascript 代码块中，不要任何解释文字。"
        )
    if lang == "go":
        if kind == "patch":
            return (
                "你是一个资深的 Go 开发工程师。根据下面的失败输出，定位问题根因，"
                "输出修复代码所需的 unified diff（diff -u 格式）：\n"
                "1. 只输出 diff 本身，放在 ```diff 围栏中，不要输出完整文件，不要任何解释文字\n"
                "2. 文件头格式：--- a/文件名 和 +++ b/文件名（patch -p1 应用）\n"
                "3. 只包含修复所需的最小改动，保留原有缩进，diff 必须格式合法\n"
                "4. 若失败源于测试期望与实现不符（如外部依赖不可用），给出合理修复（mock 外部调用或调整逻辑），"
                "不得删除接口或改变接口语义"
            )
        if kind == "test_file":
            return (
                "你是一个资深的测试工程师。根据下面的测试执行/语法错误输出和当前测试文件内容，"
                "定位问题根因并输出修复后的完整 Go 测试文件（go test ./...）。"
                "保持测试意图与覆盖范围不变，只修复导致失败的问题（语法/import/断言等）。"
                "只返回修复后的完整代码，放在 ```go 代码块中，不要任何解释文字。"
            )
        return (
            "你是一个资深的 Go 开发工程师。根据下面的构建/测试失败输出和当前入口文件代码，"
            "定位问题根因并输出修复后的完整入口文件（必须是可直接运行的 HTTP 服务，监听 0.0.0.0:8000，"
            "根路径 / 必须返回 200）。保持原有功能与接口不变，只修复导致失败的问题。"
            "只返回修复后的完整代码，放在 ```go 代码块中，不要任何解释文字。"
        )
    # python 默认
    if kind == "patch":
        return PATCH_FIX_SYSTEM
    if kind == "test_file":
        return TEST_FILE_FIX_SYSTEM
    return TEST_FIX_SYSTEM



def _load_test_cases(cfg: dict) -> str:
    """加载需求测试用例（无则空串）。"""
    if not cfg.get("requirement_id"):
        return ""
    conn = get_db()
    try:
        row = conn.execute("SELECT test_cases FROM requirements WHERE id=?", (cfg["requirement_id"],)).fetchone()
        return row["test_cases"] if row and row["test_cases"] else ""
    finally:
        conn.close()


def _build_test_prompt(lang: str, entry: str, code: str, test_cases: str) -> str:
    """构建测试文件生成提示词（超长代码截断）。"""
    sys_prompt = _TEST_PROMPTS.get(lang, _TEST_PROMPTS["python"]).replace("<entry>", entry)
    prompt = f"【需求测试用例】\n{test_cases or '（无，请基于代码接口自拟核心用例）'}\n\n【{entry}】\n{code}"
    if len(prompt) > 17000:
        prompt = prompt[:11000] + "\n# ……（代码过长已截断）……\n" + prompt[-6000:]
    return prompt


def _rewrite_test_file(lang: str, test_file: str, fixed: str, err_v: str, append) -> str | None:
    """语法错误时 LLM 重写测试文件（含错误上下文）。失败返回 None。"""
    try:
        flines = fixed.splitlines()
        em = re.search(r"L(\d+)|:(\d+):", err_v)
        err_line = int(em.group(1) or em.group(2) or 1) if em else 1
        ctx_lines = flines[max(0, err_line - 22): err_line + 18]
        ctx = "\n".join(f"{max(0, err_line - 22) + i + 1}| {line}" for i, line in enumerate(ctx_lines))
        brief = fixed[:4000] + f"\n……（共 {len(flines)} 行，中间省略）……\n" + fixed[-2000:]
        fix2 = call_llm(
            _fix_system(lang, "test_file"),
            f"【语法错误】{err_v} at line {err_line}（常见原因：输出被 token 限制截断）\n"
            f"【错误上下文】\n{ctx}\n\n【文件结构（头尾摘要）】\n{brief}\n\n"
            f"请重新输出修复后的完整 {test_file}（必须完整不截断）。",
            max_tokens=6000,
            timeout=180,
        )
        fixed2 = _extract_code_block(fix2)
        ok_v2, err_v2 = _validate_test_file(lang, fixed2)
    except Exception as e2:
        append(f"  - ⚠ 测试文件重写失败: {e2}，跳过自动化测试门禁")
        return None
    if not fixed2 or not ok_v2:
        append(f"  - ⚠ 测试文件重写仍无效（{err_v2}），跳过自动化测试门禁")
        return None
    append(f"  - 测试文件重写成功（{len(fixed)} 字节）")
    return fixed2


def _ensure_test_file(project_dir, cfg, append) -> bool:  # noqa: C901
    """按技术栈生成测试文件（python→pytest / node→node --test / go→go test），已有则复用。"""
    lang = cfg.get("language") or _detect_project_type(project_dir)["lang"]
    if lang == "docker":
        ext = (cfg.get("entry") or "main.py").rsplit(".", 1)[-1]
        lang = {"py": "python", "js": "node", "ts": "node", "go": "go"}.get(ext, "python")
    entry = cfg.get("entry") or "main.py"
    test_file = cfg.get("test_file") or {
        "python": "test_main.py",
        "node": "test_api.test.js",
        "go": "main_test.go",
    }.get(lang, "test_main.py")
    tf_path = os.path.join(project_dir, test_file)
    if os.path.exists(tf_path) and os.path.getsize(tf_path) > 50:
        return True
    try:
        with open(os.path.join(project_dir, entry), encoding="utf-8") as f:
            code = f.read()
    except Exception as e:
        append(f"  - ⚠ 读取 {entry} 失败: {e}")
        return False
    test_cases = _load_test_cases(cfg)
    prompt = _build_test_prompt(lang, entry, code, test_cases)
    try:
        out = call_llm(_TEST_PROMPTS.get(lang, _TEST_PROMPTS["python"]).replace("<entry>", entry), prompt, max_tokens=6000, timeout=180)
    except Exception as e:
        append(f"  - ⚠ LLM 生成测试文件失败: {e}（可在系统配置-模型列表中设置模型 API Key）")
        return False
    fixed = _extract_code_block(out)
    if not fixed:
        append("  - ⚠ LLM 未输出测试代码，跳过自动化测试门禁")
        return False
    ok_v, err_v = _validate_test_file(lang, fixed)
    if not ok_v:
        append(f"  - ⚠ 生成测试文件语法错误（{err_v}），LLM 重写中…")
        fixed2 = _rewrite_test_file(lang, test_file, fixed, err_v, append)
        if fixed2 is None:
            return False
        fixed = fixed2
    with open(tf_path, "w", encoding="utf-8") as f:
        f.write(fixed)
    append(f"  - 测试文件已生成 {test_file}（{len(fixed)} 字节）")
    return True

def _record_test_run(requirement_id, pipeline_id, status, summary, log_text, cases=None) -> None:
    """记录一次自动化测试执行结果（需求维度，AI 工作台可见）。记录失败不阻塞主流程。

    cases: 逐条用例结果 [{name, path, status, message}]，JSON 序列化入 cases 列。
    """
    if not requirement_id:
        return
    conn = get_db()
    try:
        conn.execute(
            "CREATE TABLE IF NOT EXISTS test_runs (id TEXT PRIMARY KEY, requirement_id TEXT, pipeline_id TEXT, "
            "status TEXT, summary TEXT, log TEXT, cases TEXT, created_at TEXT)"
        )
        # 旧库无 cases 列：安全追加（SQLite 不支持 ADD COLUMN IF NOT EXISTS）
        cols = {row["name"] for row in conn.execute("PRAGMA table_info(test_runs)").fetchall()}
        if "cases" not in cols:
            conn.execute("ALTER TABLE test_runs ADD COLUMN cases TEXT")
        conn.execute(
            "INSERT INTO test_runs (id, requirement_id, pipeline_id, status, summary, log, cases, created_at)"
            " VALUES (?,?,?,?,?,?,?,?)",
            (
                f"test_{uuid.uuid4().hex[:12]}",
                requirement_id,
                pipeline_id,
                status,
                summary or "",
                log_text or "",
                json.dumps(cases or [], ensure_ascii=False),
                datetime.now().isoformat(),
            ),
        )
        conn.commit()
    except Exception as e:
        logger.warning(f"记录测试结果失败: {e}")
    finally:
        conn.close()


def _parse_test_summary(out: str, lang: str = "python") -> str:
    """从测试输出提取简短结果摘要（pytest: 'N passed'；node:test: '# pass N'；go test: 'ok pkg'）。"""
    out = out or ""
    if lang == "node":
        parts = re.findall(r"#\s*(pass|fail)\s+(\d+)", out)
        return (
            ", ".join(f"{k} {v}" for k, v in parts)
            if parts
            else (out.splitlines()[-1][:120] if out.strip() else "无输出")
        )
    if lang == "go":
        m = re.search(r"^(ok|FAIL)\s+\S+[^\n]{0,60}", out, re.M)
        return m.group(0).strip() if m else (out.splitlines()[-1][:120] if out.strip() else "无输出")
    m = re.search(r"\d+ (?:passed|failed|error)[^\n]{0,60}", out)
    return m.group(0).strip() if m else (out.splitlines()[-1][:120] if out.strip() else "无输出")


def _parse_pytest_cases(out: str) -> list:
    """从 pytest -rA 输出解析逐条用例结果（AI 工作台按 case 展示）。

    返回 [{name, path, status, message}]；status ∈ passed/failed/error/skipped。
    SKIPPED 行格式特殊（`SKIPPED [1] file:line: 原因`，无 :: 分隔、带计数前缀）；
    未匹配到 -rA 摘要时退化解析 FAILED 行（失败项兜底，至少可见失败原因）。
    """
    out = out or ""
    cases, seen = [], set()
    # 优先只扫 short test summary 区（避免正文 traceback 行干扰）；无则全文
    body = out
    idx = out.rfind("short test summary info")
    if idx >= 0:
        body = out[idx:]
    # 分隔用 [ \t] 而非 \s：避免 \s 吞换行导致 (.*) 串行抓取下一行
    for m in re.finditer(
        r"^(PASSED|FAILED|SKIPPED|ERROR)\s+(?:\[\d+\]\s*)?(\S+)(?:[ \t]*(?::[ \t]*)?(.*))?$", body, re.M
    ):
        kind, path = m.group(1), m.group(2)
        if path in seen:
            continue
        seen.add(path)
        msg = (m.group(3) or "").strip()
        if msg.startswith("- "):
            msg = msg[2:].strip()
        cases.append(
            {
                "name": path.split("::")[-1].rstrip(":"),
                "path": path,
                "status": {"PASSED": "passed", "FAILED": "failed", "SKIPPED": "skipped", "ERROR": "error"}[kind],
                "message": msg[:500],
            }
        )
    if not cases:
        for m in re.finditer(r"^FAILED\s+(\S+?)(?:[ \t]*-[ \t]*(.*))?$", out or "", re.M):
            path = m.group(1)
            if path in seen:
                continue
            seen.add(path)
            cases.append(
                {
                    "name": path.split("::")[-1].rstrip(":"),
                    "path": path,
                    "status": "failed",
                    "message": (m.group(2) or "").strip()[:500],
                }
            )
    return cases


def _attach_case_meta(cases: list, project_dir: str) -> list:
    """解析测试文件 docstring 的 TC 编号标注，把每条执行结果对应到生成的用例。

    生成的用例文档（requirement.test_cases）带 TC-xxx 编号；AI 生成的测试函数
    docstring 首行标注 `TC-API-016: 仅获取实时天气` 与之对应。执行结果据此合并
    case_id/case_title，实现「生成的 case ↔ 执行结果」一一对应。
    """
    if not cases or not project_dir:
        return cases
    mapping = {}
    for tf in glob.glob(os.path.join(project_dir, "test_*.py")):
        try:
            with open(tf, encoding="utf-8") as f:
                code = f.read()
        except OSError:
            continue
        for m in re.finditer(r'def\s+(test_\w+)\s*\([^)]*\)\s*:[^\n]*\n\s*["\']{3}\s*([^\n]*)', code):
            fn, doc = m.group(1), m.group(2).strip()
            # 剥掉 docstring 结尾引号（单行 docstring 尾部紧贴 """）
            doc = re.sub(r"[\"']+$", "", doc).strip()
            cm = re.match(r"^(TC-[\w-]+)\s*[:：\-—]?\s*(.*)$", doc)
            if cm:
                mapping[fn] = {"case_id": cm.group(1), "case_title": cm.group(2).strip()[:120]}
    if not mapping:
        return cases
    for c in cases:
        meta = mapping.get(c.get("name") or "")
        if meta:
            c["case_id"] = meta["case_id"]
            c["case_title"] = meta["case_title"]
    return cases


def _extract_failed_functions(project_dir: str, out: str) -> list:  # noqa: C901
    """提取全部失败测试对应的 main.py 函数（支持多函数批量修复）。"""
    import ast

    main_file = os.path.join(project_dir, "main.py")
    try:
        with open(main_file, encoding="utf-8") as f:
            code = f.read()
        tree = ast.parse(code)
    except Exception:
        return []
    names = _parse_failed_test_names(out)
    if not names:
        return []
    test_info = _parse_test_info(project_dir, names)
    funcs, order = _match_route_functions(tree, code, test_info)
    return [(n, funcs[n]["start"], funcs[n]["end"], funcs[n]["segs"]) for n in order]

def _replace_function(target_path: str, new_code: str, start_line: int, end_line: int) -> None:
    """用修复后的函数源码替换文件指定行区间，并按原函数缩进适配。"""
    with open(target_path, encoding="utf-8") as f:
        lines = f.read().splitlines(keepends=True)
    orig = lines[start_line - 1]
    indent = orig[: len(orig) - len(orig.lstrip())]
    body_lines = [line for line in new_code.strip("\n").splitlines()]
    adapted = []
    if body_lines:
        base = body_lines[0][: len(body_lines[0]) - len(body_lines[0].lstrip())]
        for line in body_lines:
            if line.strip():
                adapted.append(indent + (line[len(base) :] if line.startswith(base) else line.lstrip()))
            else:
                adapted.append("")
    lines[start_line - 1 : end_line] = [line + "\n" for line in adapted]
    with open(target_path, "w", encoding="utf-8") as f:
        f.write("".join(lines))



def _setup_test_environment(test_id: str) -> dict:
    """准备测试环境。"""
    import tempfile
    tmp_dir = tempfile.mkdtemp(prefix=f"test_{test_id}_")
    return {"tmp_dir": tmp_dir, "test_id": test_id}

def _run_core_tests(test_config: dict) -> dict:
    """运行核心测试用例。"""
    results = {"passed": 0, "failed": 0, "errors": []}
    tests = test_config.get("tests", [])
    for test in tests:
        try:
            # 执行测试
            results["passed"] += 1
        except Exception as e:
            results["failed"] += 1
            results["errors"].append(str(e))
    return results

def _validate_results(test_results: dict, output_path: str) -> bool:
    """验证测试结果。"""
    if test_results["failed"] > 0:
        return False
    import os
    return os.path.exists(output_path)



def _prepare_test_context(pid, run_id, cfg):
    """准备测试门禁上下文。"""
    return {
        "pid": pid,
        "run_id": run_id,
        "service_name": cfg.get("service_name", ""),
        "slug": cfg.get("slug", ""),
        "project_dir": cfg.get("project_dir", ""),
        "lang": cfg.get("language") or "python",
        "status": "initialized"
    }

def _execute_test_step(test_context, step_name, step_data):
    """执行测试步骤。"""
    return {
        "step": step_name,
        "data": step_data,
        "status": "completed"
    }

def _finalize_test_results(results):
    """汇总测试结果。"""
    return {
        "total_steps": len(results),
        "results": results,
        "status": "completed"
    }

def _pick_fix_target(out: str, test_file: str, entry: str) -> str:
    """智能选择修复目标：测试文件问题 vs 实现缺陷。"""
    low = out.lower()
    if (
        "syntax" in low
        or "collection" in low
        or "unresolved import" in low
        or "assert 422" in low
        or "assert 400" in low
    ):
        return test_file
    if (
        "internal server error" in low
        or " 500 " in out
        or "status_code == 500" in low
        or "keyerror" in low
        or "typeerror" in low
        or "attributeerror" in low
        or "valueerror" in low
    ):
        return entry
    return test_file


def _fix_functions_batch(target_path: str, funcs: list, entry: str, diag: str, append, backup_path: str) -> int:
    """函数级批量修复：定位失败路由函数，逐个 LLM 重写。返回修复成功数。"""
    import ast
    import shutil

    fixed_count = 0
    for fname, fstart, fend, test_segs in sorted(funcs, key=lambda x: -x[1]):
        cur = open(target_path, encoding="utf-8").read()
        func_code = "\n".join(cur.splitlines()[fstart - 1 : fend])
        cur_lines = cur.splitlines()
        deco_start = fstart - 1
        while deco_start > 0 and cur_lines[deco_start - 1].lstrip().startswith("@"):
            deco_start -= 1
        deco_lines = cur_lines[deco_start : fstart - 1]
        ctx_extra = ""
        if deco_lines:
            ctx_extra += "\n\n【路由装饰器（response_model 约束，可修改）】\n" + "\n".join(deco_lines)
        if test_segs:
            ctx_extra += "\n\n【失败测试用例（必须满足其断言）】\n" + "\n---\n".join(
                s[:600] for s in test_segs[:3]
            )
        try:
            fix = call_llm(
                FUNCTION_FIX_SYSTEM,
                f"【失败输出】\n{diag}\n\n【函数 {fname}（{entry} {fstart}-{fend} 行）】\n{func_code}{ctx_extra}",
                max_tokens=4000,
                timeout=180,
            )
        except Exception:
            fix = ""
        fixed = _extract_code_block(fix)
        fstart_new = fstart
        if fixed:
            fstrip = fixed.strip()
            if fstrip.startswith("@"):
                fstart_new = deco_start + 1
            elif deco_lines and re.match(r"^(async\s+)?def ", fstrip):
                fixed = "\n".join(deco_lines) + "\n" + fixed
        if fixed and (re.match(r"^(async\s+)?def ", fixed.strip()) or fixed.strip().startswith("@")):
            try:
                ast.parse(fixed)
                shutil.copy2(target_path, backup_path)
                _replace_function(target_path, fixed, fstart_new, fend)
                ast.parse(open(target_path, encoding="utf-8").read())
                os.remove(backup_path)
                fixed_count += 1
                append(f"  - AI 函数级修复 {fname} 成功")
            except Exception as e:
                append(f"  - ⚠ 函数级修复 {fname} 校验失败: {e}，恢复该函数")
                shutil.copy2(backup_path, target_path)
        else:
            append(f"  - ⚠ 函数级修复 {fname} 未产出有效函数，跳过")
    return fixed_count


def _apply_llm_patch(lang: str, target: str, content: str, diag: str, target_path: str, backup_path: str, append, project_dir: str) -> bool:
    """unified diff 补丁修复（含一次 LLM 重试）。成功返回 True。"""
    import ast
    import shutil

    patch_text = ""
    try:
        fix = call_llm(
            _fix_system(lang, "patch"),
            f"【失败输出】\n{diag}\n\n【文件 {target} 全文】\n{content}",
            timeout=180,
        )
        m = re.search(r"```diff\s*\n(.*?)```", fix or "", re.DOTALL)
        patch_text = m.group(1).strip() if m else _extract_code_block(fix)
    except Exception:
        patch_text = ""

    def _apply(patch: str):
        return subprocess.run(
            ["patch", "-p1", "--forward", "--batch"],
            input=patch,
            capture_output=True,
            text=True,
            cwd=project_dir,
            timeout=30,
        )

    if not patch_text:
        return False
    shutil.copy2(target_path, backup_path)
    pr = _apply(patch_text)
    if pr.returncode != 0:
        append(
            f"  - ⚠ diff 补丁应用失败（{pr.stderr.strip().splitlines()[-1][:80] if pr.stderr else '未知原因'}），要求 LLM 重新生成补丁…"
        )
        try:
            fix2 = call_llm(
                _fix_system(lang, "patch"),
                f"【失败输出】\n{diag}\n\n【文件 {target} 全文】\n{content}\n\n"
                f"（你上一次的 diff 应用失败：{pr.stderr.strip()[:300]}。请重新输出 diff 块，"
                f"确保 hunk 的上下文行数与 @@ 行号一致，只输出一个合法的 unified diff）",
                timeout=180,
            )
            m2 = re.search(r"```diff\s*\n(.*?)```", fix2 or "", re.DOTALL)
            patch_text = m2.group(1).strip() if m2 else _extract_code_block(fix2)
        except Exception:
            patch_text = ""
        if patch_text:
            shutil.copy2(target_path, backup_path)
            pr = _apply(patch_text)
    if pr.returncode != 0:
        append(f"  - ⚠ diff 补丁应用失败: {pr.stderr[-200:]}，改用全量重写…")
        shutil.copy2(backup_path, target_path)
        return False
    try:
        ast.parse(open(target_path, encoding="utf-8").read())
    except Exception:
        pr = _apply(patch_text)
    if pr.returncode == 0:
        try:
            ast.parse(open(target_path, encoding="utf-8").read())
            os.remove(backup_path)
            append("  - AI 生成 diff 补丁并应用成功，重新构建并复跑测试…")
            return True
        except Exception as e:
            append(f"  - ⚠ 补丁后语法校验失败: {e}，恢复原文件并改用全量重写…")
            shutil.copy2(backup_path, target_path)
            return False
    append(f"  - ⚠ diff 补丁应用失败: {pr.stderr[-200:]}，恢复原文件并改用全量重写…")
    shutil.copy2(backup_path, target_path)
    return False



def _apply_full_rewrite(lang: str, target: str, test_file: str, content: str, diag: str, target_path: str, backup_path: str, append) -> str:
    """全量重写策略（仅小文件）：返回修复后代码，无效返回空串。"""
    import shutil

    if len(content) > 20000:
        append("  - ⚠ 文件超过 20KB，跳过全量重写（避免 LLM 输出截断破坏文件）")
        return ""
    if target == test_file:
        append("  - ⚠ 测试文件断言问题不做全量重写（避免丢失既有用例覆盖），本轮跳过")
        return ""
    brief = content if len(content) <= 15000 else content[:9000] + "\n# ……（代码过长已截断）……\n" + content[-6000:]
    try:
        fix = call_llm(
            _fix_system(lang, "test_file" if target == test_file else "main"),
            f"【构建/测试失败输出】\n{diag}\n\n【当前 {target}】\n{brief}",
            max_tokens=6000,
            timeout=180,
        )
    except Exception as e:
        append(f"  - ❌ LLM 调用失败: {e}（可在系统配置-模型中设置模型 API Key）")
        return ""
    fixed = _extract_code_block(fix)
    if fixed:
        ok_v, err_v = _validate_test_file(lang, fixed)
        if not ok_v:
            append(f"  - ⚠ 全量重写产物语法错误（{err_v}），本轮修复无效")
            return ""
    return fixed


def _verify_test_run(append, test_file: str, image_tag: str, project_dir: str, lang: str, test_cmd: list, net: str, env_flags: list, step_run) -> tuple:
    """构建测试镜像并执行测试。返回 (ok, out)。"""
    append(f"  - 构建测试镜像（含测试运行环境 + {test_file}）…")
    ok, out = step_run(["podman", "build", "-f", "Dockerfile.test", "-t", image_tag, project_dir], timeout=900)
    if not ok:
        return False, f"测试镜像构建失败: {out[-600:]}"
    append(f"  - 容器内执行测试（{lang}）…")
    exec_cmd = list(test_cmd)
    if lang == "python" and exec_cmd and "pytest" in exec_cmd[0].lower() and not any("-rA" in c for c in exec_cmd):
        exec_cmd.insert(1, "-rA")
    cmd = ["podman", "run", "--rm", "--network", net] + env_flags + [image_tag] + exec_cmd
    ok, out = step_run(cmd, timeout=300)
    return ok, out or ""


def _run_test_fix_round(out: str, lang: str, entry: str, test_file: str, project_dir: str, append) -> str:
    """AI 修复单轮：函数级 → diff 补丁 → 全量重写。返回 'continue'/'skip'/'rewritten'。"""
    import shutil

    target = _pick_fix_target(out, test_file, entry)
    target_path = os.path.join(project_dir, target)
    backup_path = target_path + ".bak"
    try:
        with open(target_path, encoding="utf-8") as f:
            content = f.read()
        diag = (out[:2500] + "\n……\n" + out[-1000:]) if len(out) > 4000 else out
        if lang == "python":
            try:
                funcs = _extract_failed_functions(project_dir, out)
            except Exception:
                funcs = []
            fixed_count = _fix_functions_batch(target_path, funcs, entry, diag, append, backup_path)
            if fixed_count:
                append(f"  - 本轮批量修复 {fixed_count} 个函数，重新构建并复跑测试…")
                return "continue"
            append("  - 函数级修复未产出任何有效修复，改用补丁/全量重写…")
        patched = _apply_llm_patch(lang, target, content, diag, target_path, backup_path, append, project_dir)
        if patched:
            return "continue"
        fixed = _apply_full_rewrite(lang, target, test_file, content, diag, target_path, backup_path, append)
        if not fixed:
            return "skip"
    except Exception as e:
        append(f"  - ❌ LLM 修复调用失败: {e}（可在系统配置-模型列表中设置模型 API Key）")
        return "skip"
    if not fixed:
        append("  - ⚠ LLM 未输出修复代码，本轮跳过")
        return "skip"
    shutil.copy2(target_path, backup_path)
    with open(target_path, "w", encoding="utf-8") as f:
        f.write(fixed)
    append(f"  - 修复代码已落盘 {target}（{len(fixed)} 字节），重新构建并复跑测试…")
    return "rewritten"

def _run_test_gate(pid, run_id, cfg, append, step_run) -> tuple:  # noqa: C901
    """自动化测试门禁：按技术栈生成测试文件 → 构建测试镜像 → 容器内执行 → 失败 AI 修复循环（≤3 轮）→ 通过后放行部署。

    python→pytest / node→node --test / go→go test；返回 (ok, summary)；ok=False 表示测试多次失败且修复未解决，
    流水线应终止；summary='skip' 表示测试不可用已跳过。每次执行结果写入 test_runs 表（AI 工作台可查看）。
    """
    name = cfg["service_name"]
    slug = _safe_slug(name)
    project_dir = cfg["project_dir"]
    image_tag = f"app-{slug}:test"
    container_name = f"sandbox-{slug}"
    lang = cfg.get("language") or _detect_project_type(project_dir)["lang"]
    entry = cfg.get("entry") or "main.py"
    test_file = cfg.get("test_file") or "test_main.py"
    test_cmd = cfg.get("test_cmd") or ["pytest", "-q", "--tb=short", "test_main.py"]
    # 0. 依赖容器（测试容器与部署容器同网络，可真实访问 Redis/MySQL）
    net, env_flags, ok, err = _prepare_dependencies(cfg, container_name, append, step_run)
    if not ok:
        return False, err
    # 1. 生成/复用测试文件
    if not _ensure_test_file(project_dir, cfg, append):
        return True, "skip"
    # 2. 写测试镜像 Dockerfile（按语言生成）
    with open(os.path.join(project_dir, "Dockerfile.test"), "w", encoding="utf-8") as f:
        f.write(_gen_dockerfile(lang, include_tests=True, entry=entry))

    # 3. 初始验证 + 失败修复循环（最多 5 轮 AI 修复）
    last_out = ""
    for round_no in range(6):
        ok, out = _verify_test_run(
            append, test_file, image_tag, project_dir, lang, test_cmd, net, env_flags, step_run
        )
        last_out = out
        cases = _attach_case_meta(_parse_pytest_cases(out), project_dir) if lang == "python" else []
        if ok:
            summary = _parse_test_summary(out, lang)
            _record_test_run(cfg.get("requirement_id"), pid, "passed", summary, out, cases)
            append(f"  - 自动化测试通过 ✓（{summary}）")
            return True, summary
        if _is_infra_error(out):
            append("  - ⚠ 检测到基础设施故障（Docker/Podman 环境问题），AI 修复无法解决，停止测试")
            _record_test_run(cfg.get("requirement_id"), pid, "failed", "基础设施故障", out[-800:])
            return False, out[-400:]
        # 失败轮次落库：每条 case 的失败原因在 AI 工作台可见（含 AI 修复过程轨迹）
        _record_test_run(
            cfg.get("requirement_id"),
            pid,
            "failed",
            f"第{round_no + 1}轮: {_parse_test_summary(out, lang)}",
            out,
            cases,
        )
        if round_no >= 3:
            break
        append(f"  - ⚠ 测试未通过（第 {round_no + 1} 次验证），AI 诊断修复中…")
        action = _run_test_fix_round(
            out, lang, entry, test_file, project_dir, append
        )
        if action == "continue":
            continue
        if action == "skip":
            break
    _record_test_run(
        cfg.get("requirement_id"),
        pid,
        "failed",
        "AI 修复轮次用尽",
        last_out[-1500:],
        _attach_case_meta(_parse_pytest_cases(last_out), project_dir) if lang == "python" else [],
    )
    return False, "自动化测试多次失败且 AI 修复未解决: " + last_out[-300:]


def _exec_deploy_pipeline(pid: str, run_id: str, cfg: dict) -> None:  # noqa: C901
    """后台执行部署流水线：构建镜像 → 启动沙箱容器 → 健康检查；失败自动进入 AI 修复循环。"""
    name = cfg["service_name"]
    slug = _safe_slug(name)
    project_dir = cfg["project_dir"]
    port = cfg["port"]
    image_tag = f"app-{slug}"
    container_name = f"sandbox-{slug}"
    log: list = []

    def append(line: str) -> None:
        log.append(line)
        _update_run_log(run_id, "\n".join(log))

    def step_run(cmd: list, timeout: int = 900) -> tuple:
        try:
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
        except subprocess.TimeoutExpired:
            return False, "命令执行超时"
        except Exception as e:
            return False, str(e)
        return (r.returncode == 0), (r.stdout or r.stderr or "").strip()

    try:
        append(f"[{datetime.now().isoformat()[:19]}] 部署流水线「{name}」开始执行（deploy 类型 · 真实构建）")
        # 阶段 1：检出代码（入口文件按技术栈推断，支持 python/node/go 工程）
        pt = _detect_project_type(project_dir)
        entry = cfg.get("entry") or pt["entry"]
        main_file = os.path.join(project_dir, entry)
        if not os.path.exists(main_file):
            raise RuntimeError(f"{entry} 不存在，无法构建")
        append(
            f"  - 检出代码: 就绪（artifacts/{name}/，技术栈 {pt['lang']}，入口 {entry}，{os.path.getsize(main_file)} 字节）"
        )
        # 阶段 1.6：自动化测试门禁（生成测试 → 容器内执行 → 失败 AI 修复 → 通过后放行部署）
        if cfg.get("auto_test", True):
            append("  - ⚡ 自动化测试门禁已开启：先执行测试，通过后再部署")
            t_ok, t_info = _run_test_gate(pid, run_id, cfg, append, step_run)
            if not t_ok:
                append(f"  - ❌ 自动化测试未通过，停止部署：{t_info[-400:]}")
                _finish_run(run_id, pid, "failed", "\n".join(log))
                return
            if t_info != "skip":
                append(f"  - 自动化测试门禁通过 ✓（{t_info}）")
        ok, info = _deploy_once(name, project_dir, port, image_tag, container_name, append, step_run, cfg)
        if not ok:
            append(f"  - ❌ {info}")
            if _is_infra_error(info):
                # 环境故障（podman 未启动/磁盘满等）：AI 改代码无法解决，直接终止避免无效修复轮次
                append("  - ⚠ 检测到基础设施故障（Docker/Podman 环境问题），AI 修复无法解决，请检查容器环境后重新部署")
                _finish_run(run_id, pid, "failed", "\n".join(log))
                return
            if cfg.get("auto_fix", True):
                append("  - ⚡ 失败自动修复已开启，进入 AI 诊断修复…")
                if _fix_rounds(pid, run_id, cfg, log, append, step_run, info, max_rounds=3):
                    _finish_run(run_id, pid, "success", "\n".join(log))
                else:
                    _finish_run(run_id, pid, "failed", "\n".join(log))
                return
            raise RuntimeError(info)
        append(f"  - 健康检查: 通过 ✓（HTTP {info}）")
        append(f"  - 部署完成 ✓ 访问地址: http://localhost:{port}")
        _register_sandbox(slug, port, project_dir, image_tag, cfg, display_name=name)
        _finish_run(run_id, pid, "success", "\n".join(log))
    except Exception as e:
        append(f"  - ❌ {e}")
        _finish_run(run_id, pid, "failed", "\n".join(log))


def _exec_auto_fix(pid: str, run_id: str, cfg: dict) -> None:
    """手动触发 AI 诊断修复：拉取现有容器日志 → 修复循环 → 重建部署。"""
    name = cfg["service_name"]
    slug = _safe_slug(name)
    container_name = f"sandbox-{slug}"
    log: list = []

    def append(line: str) -> None:
        log.append(line)
        _update_run_log(run_id, "\n".join(log))

    def step_run(cmd: list, timeout: int = 900) -> tuple:
        try:
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
        except subprocess.TimeoutExpired:
            return False, "命令执行超时"
        except Exception as e:
            return False, str(e)
        return (r.returncode == 0), (r.stdout or r.stderr or "").strip()

    try:
        append(f"[{datetime.now().isoformat()[:19]}] 手动触发 AI 诊断修复「{name}」")
        _, clogs = step_run(["podman", "logs", "--tail", "100", container_name], timeout=10)
        initial = clogs if clogs else "容器无日志输出或尚未启动"
        if _fix_rounds(pid, run_id, cfg, log, append, step_run, initial, max_rounds=3):
            _finish_run(run_id, pid, "success", "\n".join(log))
        else:
            _finish_run(run_id, pid, "failed", "\n".join(log))
    except Exception as e:
        append(f"  - ❌ {e}")
        _finish_run(run_id, pid, "failed", "\n".join(log))


@router.post("/api/deployments")
async def create_deployment(data: DeployRequest, current_user: dict = require_auth()):
    """一键部署：AI 生成的代码落盘 → 创建部署流水线 → 后台真实构建并启动沙箱容器。

    通用化：任意语言（python/node/go/已有 Dockerfile 工程均可），自动检测技术栈。
    """
    code = (data.code or "").strip()
    if not code:
        raise HTTPException(400, "代码不能为空")
    name = re.sub(r"[^\w\-]+", "-", (data.name or "").strip()).strip("-") or "app"
    if len(name) > 40:
        name = name[:40]
    pid, run_id, port = _create_deploy_pipeline(name, code, data.requirement_id, current_user["username"])
    return {"ok": True, "pipeline_id": pid, "run_id": run_id, "name": name, "port": port, "status": "running"}


_CODE_FILE_HEADER_RE = re.compile(
    r"^(?:#|//|/\*|--)?\s*([\w./\-]+\.(?:py|js|ts|jsx|tsx|go|java|json|html|css|sh|mod|sum|yml|yaml|sql|md|txt))\s*\*?/?$"
)
_NODE_BUILTIN_MODS = {
    "http",
    "https",
    "fs",
    "path",
    "crypto",
    "os",
    "url",
    "stream",
    "events",
    "util",
    "child_process",
    "process",
    "assert",
    "buffer",
    "dns",
    "net",
    "tls",
    "zlib",
    "querystring",
    "readline",
    "repl",
    "tty",
    "v8",
    "vm",
    "worker_threads",
    "perf_hooks",
    "async_hooks",
    "cluster",
    "module",
    "node:test",
    "node:http",
}


def _parse_code_files(code: str) -> dict:
    """解析 LLM 输出的多文件代码块（```lang path ... ```）为 {文件路径: 内容}。

    - 块头首行是合法文件路径（如 ```python main.py）→ 多文件工程按路径落盘，保留层级结构
    - 无路径块头/无围栏 → 按代码特征推断语言，退化为单文件工程
    """
    files: dict = {}
    text = (code or "").strip()
    if not text:
        return files
    # 支持两种块头格式：```lang（纯语言标记）与 ```lang path（语言+空格+文件路径）
    # 注意用水平空白 [^\S\n]* 而非 \s*：\s 含换行会把下一行代码吞进块头
    blocks = re.findall(r"```[a-zA-Z0-9]*[^\S\n]*([^\n]*)\n(.*?)```", text, re.DOTALL)
    for hdr, body in blocks:
        hdr = hdr.strip()
        if not hdr:
            continue
        m = _CODE_FILE_HEADER_RE.match(hdr)
        if m:
            files[m.group(1)] = body.strip()
    if files:
        return files
    # 退化：单文件输出（可能带 ```python 围栏或前置解释文字）
    bodies = [b for h, b in blocks]
    cleaned = max(bodies, key=len).strip() if bodies else text
    cleaned = re.sub(r"^```[a-zA-Z0-9]*\s*[^\n]*\n?", "", cleaned).rstrip()
    cleaned = re.sub(r"\n?```\s*$", "", cleaned).strip()
    if "package main" in cleaned and ("func main(" in cleaned or "import (" in cleaned):
        return {"main.go": cleaned}
    # 仅用 node 独有特征（require/express/app.listen），避免 FastAPI 的 app.get 误判
    if re.search(r"module\.exports|require\s*\(|express\s*\(\)|app\.listen\s*\(", cleaned):
        return {"server.js": cleaned}
    return {"main.py": cleaned}


def _ensure_project_manifests(files: dict) -> None:
    """为 node/go 工程补齐缺失的清单文件（package.json / go.mod），保证容器构建与测试可用。"""
    if "package.json" not in files and any(f.endswith((".js", ".jsx", ".ts", ".tsx")) for f in files):
        entry = next((f for f in files if f.endswith((".js", ".ts"))), "server.js")
        deps = set()
        for f, content in files.items():
            if not f.endswith((".js", ".jsx", ".ts", ".tsx")):
                continue
            for m in re.finditer(r"(?:require\s*\(\s*|from\s+)(['\"])([^'\"\s/]+)\1", content):
                dep = m.group(2).split("/")[0]
                if dep not in _NODE_BUILTIN_MODS and not dep.startswith(".") and not dep.startswith("node:"):
                    deps.add(dep)
        files["package.json"] = json.dumps(
            {
                "name": "app",
                "version": "1.0.0",
                "private": True,
                "scripts": {"start": f"node {entry}", "test": "node --test test_api.test.js"},
                "dependencies": {d: "*" for d in sorted(deps)},
            },
            indent=2,
            ensure_ascii=False,
        )
    if "go.mod" not in files and any(f.endswith(".go") for f in files):
        files["go.mod"] = "module app\n\ngo 1.22\n"


def _create_deploy_pipeline(name: str, code: str, requirement_id: str, username: str) -> tuple:
    """代码落盘 → 创建 deploy 流水线 → 后台线程真实构建部署。返回 (pid, run_id, port)。

    通用化：解析多文件代码块按工程层级落盘 → 检测技术栈（python/node/go/docker）
    → 生成或复用 Dockerfile → 按语言执行测试与部署。供「一键部署」与「一句话全自动」共用。
    """
    import shutil

    # 1. 解析多文件代码并落盘到 artifacts/<name>/（重部署时以新产物为准，清理旧文件）
    files = _parse_code_files(code)
    if not files:
        raise HTTPException(400, "未解析到有效代码，请检查生成产物")
    _ensure_project_manifests(files)
    project_dir = os.path.join(ARTIFACTS_DIR, name)
    if os.path.isdir(project_dir):
        shutil.rmtree(project_dir, ignore_errors=True)
    for rel, content in files.items():
        fp = os.path.join(project_dir, rel)
        os.makedirs(os.path.dirname(fp), exist_ok=True)
        with open(fp, "w", encoding="utf-8") as f:
            f.write(content)
    # 2. 检测技术栈：已有 Dockerfile 直接复用；否则按语言生成（容器内统一 8000 端口）
    pt = _detect_project_type(project_dir)
    if not pt["has_dockerfile"]:
        with open(os.path.join(project_dir, "Dockerfile"), "w", encoding="utf-8") as f:
            f.write(_gen_dockerfile(pt["lang"], entry=pt["entry"]))
    if pt["lang"] == "python" and "requirements.txt" not in files:
        req = "\n".join(_detect_python_deps("\n".join(files.values()))) + "\n"
        with open(os.path.join(project_dir, "requirements.txt"), "w", encoding="utf-8") as f:
            f.write(req)
    # 3. 创建部署流水线（type=deploy）
    pid = f"pipe_{uuid.uuid4().hex[:12]}"
    port = _find_free_port()
    desc = f"需求 {requirement_id} 自动部署" if requirement_id else f"{name} 沙箱部署"
    cfg = {
        "service_name": name,
        "project_dir": project_dir,
        "port": port,
        "requirement_id": requirement_id,
        "auto_fix": True,
        "auto_test": True,
        "language": pt["lang"],
        "entry": pt["entry"],
        "test_file": pt["test_file"],
        "test_cmd": pt["test_cmd"],
        "container_port": pt["container_port"],
    }
    run_id = f"run_{uuid.uuid4().hex[:12]}"
    now = datetime.now().isoformat()
    conn = get_db()
    try:
        conn.execute(
            "INSERT INTO pipelines (id, name, description, type, config, status, last_run, created_by, created_at, updated_at) "
            "VALUES (?,?,?,?,?,?,?,?,?,?)",
            (
                pid,
                f"部署 {name}",
                desc,
                "deploy",
                json.dumps(cfg, ensure_ascii=False),
                "running",
                now,
                username,
                now,
                now,
            ),
        )
        conn.execute(
            "INSERT INTO pipeline_runs (id, pipeline_id, status, log, started_at) VALUES (?,?,?,?,?)",
            (run_id, pid, "running", f"[{now[:19]}] 部署任务已创建，等待执行…", now),
        )
        conn.commit()
    finally:
        conn.close()
    # 4. 后台线程真实执行（不阻塞请求）
    threading.Thread(target=_exec_deploy_pipeline, args=(pid, run_id, cfg), daemon=True).start()
    return pid, run_id, port


# ══════════════════════════════════════════════════════════════
# 一句话全自动：PRD → 审查 → 技术方案 → 测试 → 代码 → 审查 → 部署
# ══════════════════════════════════════════════════════════════

AUTO_STAGES = [
    ("prd", "PRD 生成"),
    ("review", "PRD 审查"),
    ("td", "技术方案"),
    ("test", "测试用例"),
    ("code", "代码生成"),
    ("review_code", "代码审查"),
]


class AutoRunRequest(BaseModel):
    name: str = ""
    description: str
    language: str = "python"
    deploy: bool = True
    target_stage: str = "deploy"  # prd/review/td/test/code/review_code/deploy


class _AutoRunStopped(Exception):
    """用户手动停止信号"""


def _auto_run_update(run_id: str, **fields) -> None:
    fields["updated_at"] = datetime.now().isoformat()
    sets = ", ".join(f"{k}=?" for k in fields)
    conn = get_db()
    try:
        conn.execute(f"UPDATE auto_runs SET {sets} WHERE id=?", (*fields.values(), run_id))
        conn.commit()
    finally:
        conn.close()


def _auto_run_append(run_id: str, line: str) -> None:
    conn = get_db()
    try:
        row = conn.execute("SELECT log FROM auto_runs WHERE id=?", (run_id,)).fetchone()
        cur = (row["log"] if row else "") or ""
        # 日志上限 200 行，防止无限增长
        lines = (cur + line + "\n").split("\n")
        cur = "\n".join(lines[-200:])
        conn.execute("UPDATE auto_runs SET log=?, updated_at=? WHERE id=?", (cur, datetime.now().isoformat(), run_id))
        conn.commit()
    finally:
        conn.close()


def _auto_run_is_stopping(run_id: str) -> bool:
    conn = get_db()
    try:
        row = conn.execute("SELECT status FROM auto_runs WHERE id=?", (run_id,)).fetchone()
        return bool(row and row["status"] == "stopping")
    finally:
        conn.close()


def _auto_run_worker(  # noqa: C901
    run_id: str,
    req_id: str,
    name: str,
    description: str,
    language: str,
    deploy: bool,
    target_stage: str,
    username: str,
) -> None:
    """后台执行全自动流水线：一句话需求 → 6 阶段产物 → 自动部署。"""
    import asyncio

    from prd_engine import CODE_SYSTEM, PRD_SYSTEM, REVIEW_SYSTEM, TD_SYSTEM, TEST_SYSTEM, save_pipeline_output

    def log(line: str) -> None:
        _auto_run_append(run_id, f"[{datetime.now().isoformat()[:19]}] {line}")

    def progress(stage: str, status: str) -> None:
        _auto_run_update(run_id, current_stage=stage, stage_progress=json.dumps(progress_map, ensure_ascii=False))

    def check_stop() -> None:
        if _auto_run_is_stopping(run_id):
            raise _AutoRunStopped()

    def run_stage(stage_key: str, label: str, idx: int, prompt_fn) -> str:
        """执行单个阶段：标记进行中 → 调用 LLM → 保存产物 → 标记完成。返回产物内容。"""
        check_stop()
        progress_map[stage_key] = "running"
        progress(stage_key, "running")
        log(f"⏳ 阶段 {idx}/6：{label}…")
        content = prompt_fn()
        asyncio.run(save_pipeline_output(req_id, {"stage": stage_key, "content": content}))
        progress_map[stage_key] = "done"
        progress(stage_key, "done")
        log(f"✅ {label}完成（{len(content)} 字）")
        return content

    progress_map: dict = {}
    try:
        started = time.time()
        log(f"🎯 一句话需求：{description[:150]}")
        results = {}

        # 1. PRD 生成
        results["prd"] = run_stage("prd", "PRD 生成", 1, lambda: call_llm(PRD_SYSTEM, description, max_tokens=4000))
        if target_stage == "prd":
            return

        # 2. PRD 审查
        results["review"] = run_stage(
            "review", "PRD 审查", 2, lambda: call_llm(REVIEW_SYSTEM, results["prd"], max_tokens=4000)
        )
        if target_stage == "review":
            return

        # 3. 技术方案
        results["td"] = run_stage("td", "技术方案", 3, lambda: call_llm(TD_SYSTEM, results["prd"], max_tokens=6000))
        if target_stage == "td":
            return

        # 4. 测试用例
        results["test"] = run_stage(
            "test",
            "测试用例",
            4,
            lambda: call_llm(TEST_SYSTEM, f"PRD:\n{results['prd']}\n\n技术方案:\n{results['td']}", max_tokens=4000),
        )
        if target_stage == "test":
            return

        # 5. 代码生成（提取纯代码，供部署落盘）
        raw_code = run_stage(
            "code",
            "代码生成",
            5,
            lambda: call_llm(
                CODE_SYSTEM, f"语言: {language}\n任务类型: code\n\n技术方案:\n{results['td']}", max_tokens=8000
            ),
        )
        results["code"] = _extract_code_block(raw_code)
        if target_stage == "code":
            return

        # 6. 代码审查
        review_system = (
            f"你是一位资深的{language}代码审查专家。审查以下代码，给出改进建议，包括："
            "1.代码质量 2.潜在bug 3.性能优化 4.安全建议。"
        )
        results["review_code"] = run_stage(
            "review_code", "代码审查", 6, lambda: call_llm(review_system, results["code"])
        )
        log(f"⏱ 6 个阶段全部完成，耗时 {time.time() - started:.0f}s")
        if target_stage == "review_code":
            return

        # 7. 自动部署
        if deploy:
            check_stop()
            progress_map["deploy"] = "running"
            progress("deploy", "running")
            log("🚀 阶段 7/7：自动部署到沙箱容器…")
            safe = re.sub(r"[^\w\-]+", "-", name).strip("-") or "app"
            pid, run2, port = _create_deploy_pipeline(safe[:40], results["code"], req_id, username)
            _auto_run_update(run_id, pipeline_id=pid, port=port)
            # 等待部署流水线完成（最多 10 分钟，期间可停止）
            deadline = time.time() + 600
            deploy_status = "running"
            while time.time() < deadline:
                check_stop()
                conn = get_db()
                try:
                    row = conn.execute("SELECT status FROM pipeline_runs WHERE id=?", (run2,)).fetchone()
                finally:
                    conn.close()
                deploy_status = row["status"] if row else "unknown"
                if deploy_status != "running":
                    break
                time.sleep(3)
            progress_map["deploy"] = deploy_status
            progress("deploy", deploy_status)
            if deploy_status == "success":
                log(f"✅ 部署成功！访问地址：http://localhost:{port}")
            else:
                log(f"⚠️ 部署状态：{deploy_status}（可到流水线页面查看构建日志）")
            if deploy_status != "success":
                raise RuntimeError(f"部署未成功（{deploy_status}），请到流水线页面查看日志")
        else:
            progress_map["deploy"] = "skipped"

        _auto_run_update(
            run_id,
            status="success",
            current_stage="done",
            stage_progress=json.dumps(progress_map, ensure_ascii=False),
            finished_at=datetime.now().isoformat(),
        )
        log("🎉 一句话全自动完成！所有产物已保存到需求，可到 AI 工作台查看/微调。")
    except _AutoRunStopped:
        _auto_run_update(
            run_id,
            status="stopped",
            current_stage="stopped",
            stage_progress=json.dumps(progress_map, ensure_ascii=False),
            finished_at=datetime.now().isoformat(),
        )
        log("⏹ 已手动停止")
    except Exception as e:
        _auto_run_update(
            run_id,
            status="failed",
            error=str(e),
            current_stage="failed",
            stage_progress=json.dumps(progress_map, ensure_ascii=False),
            finished_at=datetime.now().isoformat(),
        )
        log(f"❌ 流程失败：{e}")


@router.post("/api/auto-run")
async def create_auto_run(data: AutoRunRequest, current_user: dict = require_auth()):
    """一句话全自动：创建需求 → 后台串行执行 6 阶段 → 自动部署到沙箱。"""
    desc = (data.description or "").strip()
    if not desc:
        raise HTTPException(400, "请描述你想要实现的功能")
    if data.target_stage not in [s[0] for s in AUTO_STAGES] + ["deploy"]:
        raise HTTPException(400, "无效的目标阶段")
    name = (data.name or "").strip() or desc[:30]
    req_id = f"req_{uuid.uuid4().hex[:12]}"
    run_id = f"arun_{uuid.uuid4().hex[:12]}"
    now = datetime.now().isoformat()
    conn = get_db()
    try:
        conn.execute(
            "INSERT INTO requirements (id, name, description, status, priority, creator, version, created_at, updated_at, active) "
            "VALUES (?,?,?,?,?,?,?,?,?,1)",
            (req_id, name, desc, "in_progress", "P1", current_user["username"], 1, now, now),
        )
        conn.execute(
            "INSERT INTO auto_runs (id, requirement_id, name, language, status, current_stage, stage_progress, log, created_by, created_at, updated_at) "
            "VALUES (?,?,?,?,?,?,?,?,?,?,?)",
            (
                run_id,
                req_id,
                name,
                data.language,
                "running",
                "prd",
                "{}",
                f"[{now[:19]}] 🚀 一句话全自动流水线已启动",
                current_user["username"],
                now,
                now,
            ),
        )
        conn.commit()
    finally:
        conn.close()
    threading.Thread(
        target=_auto_run_worker,
        daemon=True,
        args=(run_id, req_id, name, desc, data.language, data.deploy, data.target_stage, current_user["username"]),
    ).start()
    return {"ok": True, "run_id": run_id, "requirement_id": req_id, "status": "running"}


@router.get("/api/auto-runs/{run_id}")
async def get_auto_run(run_id: str, current_user: dict = require_auth()):
    """查询单条全自动运行进度"""
    conn = get_db()
    try:
        row = conn.execute("SELECT * FROM auto_runs WHERE id=?", (run_id,)).fetchone()
    finally:
        conn.close()
    if not row:
        raise HTTPException(404, "未找到运行记录")
    d = dict(row)
    d["stage_progress"] = json.loads(d.get("stage_progress") or "{}")
    return d


@router.post("/api/auto-runs/{run_id}/stop")
async def stop_auto_run(run_id: str, current_user: dict = require_auth()):
    """停止全自动运行（阶段间生效）"""
    conn = get_db()
    try:
        row = conn.execute("SELECT * FROM auto_runs WHERE id=?", (run_id,)).fetchone()
        if not row:
            raise HTTPException(404, "未找到运行记录")
        if row["status"] == "running":
            conn.execute(
                "UPDATE auto_runs SET status='stopping', updated_at=? WHERE id=?", (datetime.now().isoformat(), run_id)
            )
            conn.commit()
    finally:
        conn.close()
    return {"ok": True, "status": "stopping"}


@router.post("/api/code/generate")
def generate_code(data: CodeGenRequest, current_user: dict = require_auth()):
    """AI 代码生成"""
    try:
        system_prompt = f"""你是一位资深{data.language}开发工程师，拥有10年+生产级项目经验。

代码生成要求：
1. 生产级质量：包含错误处理、输入验证、边界条件处理
2. 可维护性：清晰的变量命名、适当的注释、单一职责原则
3. 安全性：防止注入攻击、敏感信息不硬编码、使用参数化查询
4. 性能意识：避免不必要的循环、合理使用缓存、注意内存管理
5. 测试友好：函数粒度适中、依赖可注入、纯函数优先

输出格式：
```{data.language}
// 完整可运行代码
```
后面附简要说明（3-5行）：核心技术选型、时间复杂度、使用方式

只返回代码和说明，不要冗余解释。"""
        result = call_llm(system_prompt, data.prompt)

        gen_id = f"cg_{uuid.uuid4().hex[:12]}"
        # 独立连接：避免 LLM 调用期间其他函数关闭线程复用连接
        from common.db import get_db_context

        with get_db_context() as conn:
            conn.execute(
                "INSERT INTO code_generations (id, language, prompt, result, model) VALUES (?,?,?,?,?)",
                (gen_id, data.language, data.prompt, result, data.model),
            )
        return {"ok": True, "id": gen_id, "result": result}
    except Exception as e:
        raise HTTPException(500, "操作失败，请稍后重试") from e


@router.get("/api/code/generations")
async def list_code_generations(current_user: dict = require_auth()):
    """获取代码生成历史"""
    conn = get_db()
    try:
        items = []
        for row in conn.execute("SELECT * FROM code_generations ORDER BY created_at DESC LIMIT 50").fetchall():
            items.append(dict(row))
        return items
    finally:
        conn.close()


@router.post("/api/code/review")
async def review_code(data: CodeReviewRequest, current_user: dict = require_auth()):
    """AI 代码审查（v12.0：stream: true 走 SSE 流式）"""
    try:
        system_prompt = f"""你是资深{data.language}代码审查专家，负责把关生产级代码质量。

审查维度（按严重程度排列）：
🔴 严重（必须修复）：
1. 安全漏洞：注入攻击、越权访问、敏感信息泄露、不安全的反序列化
2. 逻辑错误：边界条件错误、空指针/None引用、竞态条件、死锁风险
3. 数据问题：数据丢失风险、事务不一致、类型转换错误

🟡 重要（应该修复）：
4. 性能瓶颈：N+1查询、不必要的循环、内存泄漏、阻塞IO
5. 可维护性：过长函数、过深嵌套、重复代码、魔法数字

🟢 建议（可选优化）：
6. 代码风格：命名规范、注释质量、代码组织
7. 最佳实践：设计模式运用、库函数使用建议

输出格式：
## 审查总结（1-2句整体评价）
## 问题清单
### 🔴 [严重] 问题标题
- 位置：第X行 / 函数名
- 问题：具体描述
- 风险：可能导致什么后果
- 修复：给出修改代码

（每个问题独立一节，按严重程度排序，最多列出10个问题）"""
        if data.stream:
            return stream_llm_response(system_prompt, data.code, 4000, "code_review")

        result = await call_llm_async(system_prompt, data.code)

        review_id = f"cr_{uuid.uuid4().hex[:12]}"
        # 独立连接：避免 LLM 调用期间其他函数关闭线程复用连接
        from common.db import get_db_context

        with get_db_context() as conn:
            conn.execute(
                "INSERT INTO code_reviews (id, language, code, result, model) VALUES (?,?,?,?,?)",
                (review_id, data.language, data.code, result, data.model),
            )
        return {"ok": True, "id": review_id, "result": result}
    except Exception as e:
        logger.error(f"[review_code] {traceback.format_exc()}")
        raise HTTPException(500, "操作失败，请稍后重试") from e


@router.post("/api/code/improve")
def improve_code(data: CodeImproveRequest, current_user: dict = require_auth()):
    """AI 根据代码审查意见修改代码：审查结果 → 修改后的完整代码。"""
    try:
        system_prompt = (
            f"你是一位资深{data.language}开发工程师，负责根据审查意见修复代码。\n\n"
            "修复原则：\n"
            "1. 逐一处理所有审查意见中标记为🔴严重和🟡重要的问题\n"
            "2. 修复后代码必须保持原有功能不变，不引入新bug\n"
            "3. 每次修改标注修复了哪个问题（行内注释 // fix: ...）\n"
            "4. 如某个建议因架构限制无法采纳，标注说明原因\n\n"
            "输出格式：\n"
            "## 修改后的完整代码\n"
            "```{data.language}\n...\n```\n"
            "## 修改说明（表格格式）\n"
            "| 问题 | 修复方式 | 影响范围 |\n"
            "|------|----------|----------|\n"
            "| ... | ... | ... |"
        )
        prompt = f"【原始代码】\n{data.code}\n\n【代码审查意见】\n{data.review}"
        result = call_llm(system_prompt, prompt)
        return {"ok": True, "result": result}
    except Exception as e:
        logger.error(f"[code_improve] {traceback.format_exc()}")
        raise HTTPException(500, "操作失败，请稍后重试") from e


@router.get("/api/code/reviews")
async def list_code_reviews(current_user: dict = require_auth()):
    """获取代码审查历史"""
    conn = get_db()
    try:
        items = []
        for row in conn.execute("SELECT * FROM code_reviews ORDER BY created_at DESC LIMIT 50").fetchall():
            items.append(dict(row))
        return items
    finally:
        conn.close()


# Pipeline CRUD
@router.get("/api/pipelines")
async def list_pipelines(current_user: dict = require_auth()):
    conn = get_db()
    try:
        items = []
        for row in conn.execute("SELECT * FROM pipelines WHERE active=1 ORDER BY created_at DESC").fetchall():
            p = dict(row)
            p["config"] = json.loads(p.get("config", "{}"))
            # deploy 流水线：补充关联需求名称（前端卡片展示部署了什么）
            if p["type"] == "deploy" and p["config"].get("requirement_id"):
                req = conn.execute(
                    "SELECT name FROM requirements WHERE id=? AND active=1", (p["config"]["requirement_id"],)
                ).fetchone()
                if req:
                    p["config"]["requirement_name"] = req["name"]
            # 最近一次运行摘要
            run = conn.execute(
                "SELECT status, started_at, finished_at FROM pipeline_runs WHERE pipeline_id=? ORDER BY started_at DESC LIMIT 1",
                (p["id"],),
            ).fetchone()
            p["last_run"] = dict(run) if run else None
            items.append(p)
        return items
    finally:
        conn.close()


@router.post("/api/pipelines")
async def create_pipeline(data: PipelineCreate, current_user: dict = require_auth()):
    conn = get_db()
    try:
        pid = f"pipe_{uuid.uuid4().hex[:12]}"
        conn.execute(
            "INSERT INTO pipelines (id, name, description, type, config, created_by) VALUES (?,?,?,?,?,?)",
            (pid, data.name, data.description, data.type, json.dumps(data.config), current_user["username"]),
        )
        conn.commit()
        return {"ok": True, "id": pid}
    finally:
        conn.close()


class PipelineUpdate(BaseModel):
    name: str = ""
    description: str = ""
    type: str = ""
    config: dict = None


@router.put("/api/pipelines/{pid}")
async def update_pipeline(pid: str, data: PipelineUpdate, current_user: dict = require_auth()):
    conn = get_db()
    try:
        row = conn.execute("SELECT id FROM pipelines WHERE id=? AND active=1", (pid,)).fetchone()
        if not row:
            raise HTTPException(404, "流水线不存在")
        updates, values = [], []
        if data.name:
            updates.append("name=?")
            values.append(data.name)
        if data.description is not None:
            updates.append("description=?")
            values.append(data.description)
        if data.type:
            updates.append("type=?")
            values.append(data.type)
        if data.config is not None:
            updates.append("config=?")
            values.append(json.dumps(data.config, ensure_ascii=False))
        if not updates:
            raise HTTPException(400, "没有需要更新的字段")
        updates.append("updated_at=?")
        values.append(datetime.now().isoformat())
        values.append(pid)
        conn.execute(f"UPDATE pipelines SET {','.join(updates)} WHERE id=?", values)
        conn.commit()
        return {"ok": True, "id": pid}
    finally:
        conn.close()


@router.post("/api/pipelines/{pid}/run")
def run_pipeline(pid: str, current_user: dict = require_auth()):
    """执行流水线：deploy 类型真实构建并部署沙箱；其余类型按模拟日志执行。"""
    conn = get_db()
    try:
        row = conn.execute("SELECT * FROM pipelines WHERE id=? AND active=1", (pid,)).fetchone()
        if not row:
            raise HTTPException(404, "流水线不存在")
        p = dict(row)
        cfg = json.loads(p.get("config", "{}") or "{}")
        run_id = f"run_{uuid.uuid4().hex[:12]}"
        started = datetime.now().isoformat()
        ptype = p.get("type") or "ci"

        # deploy 类型：后台线程真实执行（构建镜像 + 启动沙箱容器）
        if ptype == "deploy":
            if not cfg.get("service_name") or not cfg.get("project_dir"):
                # 手动新建的 deploy 流水线缺少部署配置，直接标记失败
                err_log = f"[{started[:19]}] 部署配置缺失（service_name/project_dir）\n  - ❌ 该流水线由 AI 工作台「一键部署」自动创建，不支持手动新建运行"
                conn.execute(
                    "INSERT INTO pipeline_runs (id, pipeline_id, status, log, started_at, finished_at) VALUES (?,?,?,?,?,?)",
                    (run_id, pid, "failed", err_log, started, started),
                )
                conn.execute("UPDATE pipelines SET status='failed', last_run=? WHERE id=?", (started, pid))
                conn.commit()
                return {
                    "ok": False,
                    "id": run_id,
                    "status": "failed",
                    "error": "部署配置缺失（service_name/project_dir）",
                }
            conn.execute(
                "INSERT INTO pipeline_runs (id, pipeline_id, status, log, started_at) VALUES (?,?,?,?,?)",
                (run_id, pid, "running", f"[{started[:19]}] 部署任务已创建，等待执行…", started),
            )
            conn.execute("UPDATE pipelines SET status='running', last_run=? WHERE id=?", (started, pid))
            conn.commit()
            threading.Thread(target=_exec_deploy_pipeline, args=(pid, run_id, cfg), daemon=True).start()
            return {"ok": True, "id": run_id, "status": "running", "started_at": started}

        # 其余类型：模拟执行阶段（按类型生成日志）
        stages = {
            "ci": [
                ("检出代码", f"git clone --depth 1 {cfg.get('repo', 'https://example.com/repo.git')}"),
                ("安装依赖", "pip install -r requirements.txt（模拟）"),
                ("静态检查", "ruff check . --select E,F（模拟）"),
                ("单元测试", f"pytest -q {cfg.get('test_path', 'tests/')}（模拟）"),
                ("构建产物", "构建完成，产物打包成功（模拟）"),
            ],
            "cd": [
                ("拉取产物", "docker pull registry.example.com/app:latest（模拟）"),
                ("滚动发布", "kubectl rollout status deploy/app（模拟）"),
                ("健康检查", "GET /healthz -> 200 OK（模拟）"),
                ("发布完成", "新版本 v1.0.0 已上线（模拟）"),
            ],
            "test": [
                ("收集用例", "pytest --collect-only（模拟）"),
                ("执行用例", f"pytest -q {cfg.get('test_path', 'tests/')}（模拟）"),
                ("覆盖率", "coverage report -m（模拟）"),
            ],
            "build": [
                ("编译", "编译源码（模拟）"),
                ("打包", "构建 Docker 镜像 app:latest（模拟）"),
                ("推送镜像", "push registry.example.com/app:latest（模拟）"),
            ],
        }
        lines = [f"[{started[:19]}] 流水线「{p['name']}」开始执行（{ptype} 类型）"]
        for name, cmd in stages.get(ptype, stages["ci"]):
            lines.append(f"  - {name}: {cmd}")
        # 模拟耗时
        time.sleep(0.8)
        lines.append("  - 全部阶段通过 ✓")
        log = "\n".join(lines)
        finished = datetime.now().isoformat()
        conn.execute(
            "INSERT INTO pipeline_runs (id, pipeline_id, status, log, started_at, finished_at) VALUES (?,?,?,?,?,?)",
            (run_id, pid, "success", log, started, finished),
        )
        conn.execute("UPDATE pipelines SET status='success', last_run=? WHERE id=?", (finished, pid))
        conn.commit()
        return {
            "ok": True,
            "id": run_id,
            "status": "success",
            "log": log,
            "started_at": started,
            "finished_at": finished,
        }
    finally:
        conn.close()


@router.post("/api/pipelines/{pid}/auto-fix")
async def auto_fix_pipeline(pid: str, current_user: dict = require_auth()):
    """AI 诊断修复：拉取容器日志 → LLM 分析根因并修改代码 → 重建 → 重启 → 健康检查。

    用户可控：部署失败后可手动触发；deploy 流水线失败且 auto_fix 开启时自动触发。
    """
    conn = get_db()
    try:
        row = conn.execute("SELECT * FROM pipelines WHERE id=? AND active=1", (pid,)).fetchone()
        if not row:
            raise HTTPException(404, "流水线不存在")
        p = dict(row)
        if p.get("type") != "deploy":
            raise HTTPException(400, "仅支持对沙箱部署（deploy）流水线执行 AI 修复")
        cfg = json.loads(p.get("config", "{}") or "{}")
        if not cfg.get("service_name") or not cfg.get("project_dir"):
            raise HTTPException(400, "部署配置缺失，无法修复")
        # 确保 auto_fix 生效（修复后也保持开启）
        cfg["auto_fix"] = True
        conn.execute(
            "UPDATE pipelines SET config=?, status='running', last_run=? WHERE id=?",
            (json.dumps(cfg, ensure_ascii=False), datetime.now().isoformat(), pid),
        )
        run_id = f"run_{uuid.uuid4().hex[:12]}"
        started = datetime.now().isoformat()
        conn.execute(
            "INSERT INTO pipeline_runs (id, pipeline_id, status, log, started_at) VALUES (?,?,?,?,?)",
            (run_id, pid, "running", f"[{started[:19]}] AI 诊断修复任务已创建，等待执行…", started),
        )
        conn.commit()
        threading.Thread(target=_exec_auto_fix, args=(pid, run_id, cfg), daemon=True).start()
        return {"ok": True, "id": run_id, "status": "running", "started_at": started}
    finally:
        conn.close()


@router.post("/api/sandbox/projects/{project_id}/logs/analyze")
def analyze_sandbox_logs(project_id: str, current_user: dict = require_auth()):
    """AI 分析沙箱容器日志：拉取日志 → LLM 定位问题根因 → 返回诊断报告。"""
    import subprocess as _sp

    logs = ""
    if project_id.startswith("deploy-"):
        container = f"sandbox-{project_id[len('deploy-') :]}"
        r = _sp.run(["podman", "logs", "--tail", "200", container], capture_output=True, text=True, timeout=15)
        logs = r.stdout if r.returncode == 0 else (r.stderr or "")
    else:
        try:
            from sandbox import process_manager
        except Exception:  # noqa: BLE001  sandbox 为研发/容器功能模块，当前版本未内置
            raise HTTPException(400, "沙箱（容器）功能未在当前版本提供")

        logs = "\n".join(process_manager.get_logs(project_id, tail=200))
    if not logs.strip():
        return {"ok": True, "analysis": "容器暂无日志输出，可能尚未启动或无错误信息。", "logs": ""}
    sys_prompt = (
        "你是一个资深的 SRE 运维专家。分析下面的容器运行日志，定位问题根因，"
        "给出：1.问题现象 2.根本原因 3.修复建议（具体到代码/配置层面）。"
        "简洁清晰，使用中文，用 markdown 列表组织。不要猜测没有依据的问题。"
    )
    try:
        analysis = call_llm(sys_prompt, f"【容器日志】\n{logs[-6000:]}")
    except Exception as e:
        raise HTTPException(500, "操作失败，请稍后重试") from e
    return {"ok": True, "analysis": analysis, "logs": logs}


@router.get("/api/pipelines/{pid}/runs")
async def list_pipeline_runs(pid: str, current_user: dict = require_auth()):
    """获取流水线运行历史。"""
    conn = get_db()
    try:
        items = []
        for row in conn.execute(
            "SELECT * FROM pipeline_runs WHERE pipeline_id=? ORDER BY started_at DESC LIMIT 20", (pid,)
        ).fetchall():
            items.append(dict(row))
        return items
    finally:
        conn.close()


@router.delete("/api/pipelines/{pid}")
async def delete_pipeline(pid: str, current_user: dict = require_auth()):
    conn = get_db()
    try:
        conn.execute("UPDATE pipelines SET active=0 WHERE id=?", (pid,))
        conn.commit()
        return {"ok": True}
    finally:
        conn.close()


# ══════════════════════════════════════════════════════════════
# Phase 3: 内容创作增强
# ══════════════════════════════════════════════════════════════


# ── 商业化基线：历史表用户隔离（幂等迁移，兼容存量库）──
# 新库由 common/db.py 建表 SQL 直接带 user_id；旧库在 import 时补列。
# PRAGMA 查询不锁表，ALTER 冲突（并发竞态）时另一进程已加列，忽略即可。
def _ensure_content_user_columns() -> None:
    with get_db_context() as conn:
        for table in ("copywriting_tasks", "translations", "ppt_generations", "excel_operations"):
            cols = [r[1] for r in conn.execute(f"PRAGMA table_info({table})").fetchall()]
            if "user_id" not in cols:
                try:
                    conn.execute(f"ALTER TABLE {table} ADD COLUMN user_id TEXT DEFAULT ''")
                except Exception:
                    pass
        # PPT 另需 file_path 列（商业化：PPTX 文件下载）
        cols = [r[1] for r in conn.execute("PRAGMA table_info(ppt_generations)").fetchall()]
        if "file_path" not in cols:
            try:
                conn.execute("ALTER TABLE ppt_generations ADD COLUMN file_path TEXT DEFAULT ''")
            except Exception:
                pass
        conn.commit()


def _ensure_content_user_columns_safe() -> None:
    """兼容表不存在场景（测试库尚未 init_db 时模块级调用）"""
    try:
        _ensure_content_user_columns()
    except Exception:
        pass


_ensure_content_user_columns_safe()


# ── PPT 商业化：真实 PPTX 文件输出 ──
PPTX_DIR = os.path.join(os.path.dirname(__file__), "uploads", "ppt")
os.makedirs(PPTX_DIR, exist_ok=True)


def _user_scope_clause(conn, current_user: dict) -> tuple[str, list]:
    """历史记录用户隔离：admin 全量；普通用户返回 ' AND user_id=?' 拼接到既有条件后
    （存量无归属数据不展示）。"""
    role = current_user.get("role", "")
    if role in ("admin", "super_admin"):
        return "", []
    return " AND user_id=?", [str(current_user.get("user_id", ""))]


class CopywritingRequest(BaseModel):
    type: str = "marketing"
    title: str = ""
    prompt: str
    model: str = ""
    platform: str = ""  # v15：平台适配（wechat/xiaohongshu/douyin/zhihu/weibo/toutiao）


class TranslationRequest(BaseModel):
    source_lang: str = "中文"
    target_lang: str = "English"
    text: str
    model: str = ""
    use_glossary: bool = True


# ── 文案工厂：专家级文案类型规格（异步任务 worker 复用）──
_COPYWRITING_SPECS = {
    "marketing": {
        "role": "资深营销文案专家（8年+4A公司经验）",
        "focus": ["卖点提炼", "行动号召(CTA)", "情感共鸣", "信任背书"],
        "format": "标题（3-5个备选）+ 正文（300-500字）+ 标语Slogan + 发布建议",
    },
    "social": {
        "role": "社交媒体内容策划专家",
        "focus": ["话题性", "互动引导", "平台适配", "时效热点"],
        "format": "标题 + 正文（150-300字）+ 话题标签# + 配图建议 + 发布时间建议",
    },
    "seo": {
        "role": "SEO内容策略专家",
        "focus": ["关键词密度", "标题标签优化", "元描述", "内部链接建议"],
        "format": "SEO标题（含主关键词）+ 元描述（150-160字符）+ 正文（800-1500字）+ 关键词列表 + 结构建议",
    },
    "email": {
        "role": "邮件营销转化率优化专家",
        "focus": ["打开率优化", "点击率优化", "个性化", "A/B测试建议"],
        "format": "邮件主题行（3个备选）+ 预览文本 + 正文HTML（含CTA按钮）+ 发送时间建议",
    },
    "ad": {
        "role": "创意广告策略专家",
        "focus": ["差异化定位", "用户痛点", "场景植入", "记忆点设计"],
        "format": "广告语（5-8个备选）+ 长文案（200-400字）+ 视觉创意brief + 投放渠道建议",
    },
}


# ── 文案工厂：平台适配风格（v15：公众号/小红书/抖音/知乎/微博/头条）──
_PLATFORM_STYLES = {
    "wechat": {
        "label": "微信公众号",
        "title": "标题突出利益点或悬念钩子，12-20字，避免标题党",
        "rules": [
            "正文开篇1-2句点明主题与读者收益，避免冗长铺垫",
            "善用小标题分段（每300字左右一个），增强可扫读性",
            "段落不超过4行，结尾设置『点赞/在看/转发』引导与话题互动",
            "金句收尾，方便转载与二次传播",
        ],
    },
    "xiaohongshu": {
        "label": "小红书",
        "title": "标题含关键词+数字+情绪词（如：3个方法/必看/绝了），可带emoji",
        "rules": [
            "正文开头直接给结论或场景代入，第一行即抓住注意力",
            "口语化表达，多用『我』『姐妹』等亲切人称",
            "分段清晰，每段2-3行，重点用emoji/符号标注",
            "结尾带3-5个精准话题标签 #",
        ],
    },
    "douyin": {
        "label": "抖音",
        "title": "视频文案前3秒钩子：悬念/冲突/利益，字数不超过20字",
        "rules": [
            "口播文案短句为主，单句不超过15字，节奏明快",
            "开头3秒必须抛出钩子（反问/痛点/反常识）",
            "中段干货密集，每15-20秒一个记忆点",
            "结尾引导关注/评论/转发，可加互动问题",
        ],
    },
    "zhihu": {
        "label": "知乎",
        "title": "标题即问题，正文以直接回答开头，体现专业性与真诚",
        "rules": [
            "结构：结论先行 → 分点论证（数据/案例）→ 总结",
            "使用专业术语但辅以通俗解释，建立可信度",
            "适度使用小标题与加粗突出关键结论",
            "结尾可附『如果对你有帮助，欢迎点赞收藏』",
        ],
    },
    "weibo": {
        "label": "微博",
        "title": "话题词+核心信息前置，字数不超过30字，可带#话题#",
        "rules": [
            "正文短平快，100-200字以内",
            "内容要具备传播性：共鸣/争议/实用/趣味任一要素",
            "善用@提及与#话题#提升曝光",
            "结尾引导转发互动",
        ],
    },
    "toutiao": {
        "label": "今日头条",
        "title": "标题信息明确+数字量化（如：5个技巧/3年经验），不超过25字",
        "rules": [
            "开篇2-3句交代背景与核心信息，符合信息流阅读习惯",
            "内容密度高，避免注水，每段都要有信息量",
            "小标题分段，适合快速扫读",
            "结尾总结要点+引导关注",
        ],
    },
}


def _build_copywriting_prompt(copy_type: str, user_prompt: str, platform: str = "") -> str:
    spec = _COPYWRITING_SPECS.get(copy_type, _COPYWRITING_SPECS["marketing"])
    platform_block = ""
    if platform and platform in _PLATFORM_STYLES:
        p = _PLATFORM_STYLES[platform]
        platform_block = f"""

## 平台适配（发布目标：{p["label"]}）
标题要求：{p["title"]}
正文规则：
{chr(10).join(f"- {r}" for r in p["rules"])}"""
    return f"""你是{spec["role"]}。

核心能力：
{chr(10).join(f"- {f}" for f in spec["focus"])}

输出格式（必须严格遵循）：
{spec["format"]}

创作原则：
1. 始终以目标用户视角思考，回答"这对我有什么用？"
2. 语言简洁有力，避免行业黑话和空洞形容词
3. 每个观点用数据或场景支撑，拒绝泛泛而谈
4. 结尾始终包含可衡量的下一步行动建议
5. 如涉及品牌，需注明品牌调性建议（如：年轻活泼/专业稳重/温暖亲切）
{platform_block}"""


async def _copywriting_worker(payload: dict, progress: Callable | None = None) -> dict:
    """文案生成 worker：LLM 创作 → 记录入库（带用户归属）。"""

    def _report(pct: float, stage: str) -> None:
        _notify_progress(progress, pct, stage)

    _report(5, "解析需求")
    system_prompt = _build_copywriting_prompt(
        payload.get("type", "marketing"), payload.get("prompt", ""), payload.get("platform", "")
    )
    _report(30, "AI 创作中")
    result = call_llm(system_prompt, payload.get("prompt", ""))
    _report(85, "保存记录")
    task_id = f"copy_{uuid.uuid4().hex[:12]}"
    with get_db_context() as conn:
        conn.execute(
            "INSERT INTO copywriting_tasks (id, type, title, prompt, result, model, user_id) VALUES (?,?,?,?,?,?,?)",
            (
                task_id,
                payload.get("type", "marketing"),
                payload.get("title", ""),
                payload.get("prompt", ""),
                result,
                payload.get("model", ""),
                payload.get("user_id", ""),
            ),
        )
    _report(100, "完成")
    return {"id": task_id, "result": result}


async def _copywriting_handler(task_id: str, payload: dict, update: Callable, ctx: dict) -> dict:
    """异步任务处理器：包装文案生成，回报进度。"""
    return await _copywriting_worker(payload, progress=update)


@router.post("/api/copywriting/generate")
async def generate_copywriting(data: CopywritingRequest, current_user: dict = require_auth()):
    """AI 文案生成（异步任务：进度跟踪 / 失败自动重试 / 并发控制）"""
    if not data.prompt.strip():
        raise HTTPException(400, "文案需求不能为空")
    payload = {
        "type": data.type,
        "title": data.title,
        "prompt": data.prompt,
        "model": data.model,
        "platform": data.platform,
        "user_id": str(current_user.get("user_id", "")),
        "username": current_user.get("username", ""),
    }
    task = create_task(
        "copywriting_generate",
        payload,
        username=current_user.get("username", ""),
        user_id=str(current_user.get("user_id", "")),
        role=current_user.get("role", ""),
    )
    return {"ok": True, "task_id": task["id"], "status": task["status"]}


@router.get("/api/copywriting/history")
async def list_copywriting_history(current_user: dict = require_auth()):
    conn = get_db()
    try:
        where, args = _user_scope_clause(conn, current_user)
        items = []
        for row in conn.execute(
            f"SELECT * FROM copywriting_tasks WHERE 1=1{where} ORDER BY created_at DESC LIMIT 50", args
        ).fetchall():
            items.append(dict(row))
        return items
    finally:
        conn.close()


@router.delete("/api/copywriting/{task_id}")
async def delete_copywriting(task_id: str, current_user: dict = require_auth()):
    conn = get_db()
    try:
        where, args = _user_scope_clause(conn, current_user)
        conn.execute(f"DELETE FROM copywriting_tasks WHERE id=?{where}", [task_id] + args)
        conn.commit()
        return {"ok": True}
    finally:
        conn.close()


def _build_translation_prompt(
    source_lang: str, target_lang: str, glossary_items: list | None = None
) -> str:
    """翻译 System Prompt；glossary_items 非空时附加强制术语表规则（v15）。"""
    glossary_block = ""
    if glossary_items:
        lines = "\n".join(
            f"- {g['source_term']} → {g['target_term']}" for g in glossary_items[:50]
        )
        glossary_block = f"""

## 强制术语表（最高优先级规则）
以下术语是用户自定义的固定译文，原文中出现时必须使用指定译文，不得另行翻译、不得意译、不得保留原文：
{lines}
"""
    return f"""你是专业翻译，精通{source_lang}和{target_lang}双语互译。

翻译原则：
1. 忠实原文：准确传达原文信息，不添加、不删减、不曲解
2. 语言地道：目标语言表达自然流畅，符合母语者表达习惯
3. 风格一致：保持原文的正式/非正式语体、行业术语、修辞手法
4. 文化适配：涉及文化特定概念时，提供等效表达并标注注释
5. 格式保留：保持原文的段落结构、列表、编号、Markdown标记
{glossary_block}

输出要求：
- 只返回翻译结果，不要任何前置说明或后记
- 如果遇到专有名词、品牌名、人名，保持原文不翻译
- 如果原文有歧义，选择最合理的解释并标注 [注: ...]
- 技术术语统一使用目标语言行业标准译法"""


def _load_glossary(user_id: str) -> list:
    """加载用户术语表（个人数据，按 user_id 精确过滤）。"""
    if not user_id:
        return []
    conn = get_db()
    try:
        rows = conn.execute(
            "SELECT id, source_term, target_term FROM translation_glossary WHERE user_id = ? ORDER BY created_at DESC",
            (user_id,),
        ).fetchall()
        return [dict(r) for r in rows]
    finally:
        conn.close()


async def _translation_worker(payload: dict, progress: Callable | None = None) -> dict:
    """翻译 worker：LLM 翻译（含强制术语表）→ 记录入库（带用户归属）。"""

    def _report(pct: float, stage: str) -> None:
        _notify_progress(progress, pct, stage)

    _report(10, "AI 翻译中")
    glossary = []
    if payload.get("use_glossary", True):
        glossary = _load_glossary(payload.get("user_id", ""))
    system_prompt = _build_translation_prompt(
        payload.get("source_lang", "中文"), payload.get("target_lang", "English"), glossary
    )
    result = call_llm(system_prompt, payload.get("text", ""))
    _report(85, "保存记录")
    trans_id = f"trans_{uuid.uuid4().hex[:12]}"
    with get_db_context() as conn:
        conn.execute(
            "INSERT INTO translations (id, source_lang, target_lang, source_text, result, model, user_id) VALUES (?,?,?,?,?,?,?)",
            (
                trans_id,
                payload.get("source_lang", "中文"),
                payload.get("target_lang", "English"),
                payload.get("text", ""),
                result,
                payload.get("model", ""),
                payload.get("user_id", ""),
            ),
        )
    _report(100, "完成")
    return {"id": trans_id, "result": result}


async def _translation_handler(task_id: str, payload: dict, update: Callable, ctx: dict) -> dict:
    """异步任务处理器：包装翻译，回报进度。"""
    return await _translation_worker(payload, progress=update)


@router.post("/api/translation/translate")
async def translate_text(data: TranslationRequest, current_user: dict = require_auth()):
    """AI 翻译（异步任务：进度跟踪 / 失败自动重试 / 并发控制）"""
    if not data.text.strip():
        raise HTTPException(400, "待翻译文本不能为空")
    payload = {
        "source_lang": data.source_lang,
        "target_lang": data.target_lang,
        "text": data.text,
        "model": data.model,
        "use_glossary": data.use_glossary,
        "user_id": str(current_user.get("user_id", "")),
        "username": current_user.get("username", ""),
    }
    task = create_task(
        "translation_translate",
        payload,
        username=current_user.get("username", ""),
        user_id=str(current_user.get("user_id", "")),
        role=current_user.get("role", ""),
    )
    return {"ok": True, "task_id": task["id"], "status": task["status"]}


@router.get("/api/translation/history")
async def list_translation_history(current_user: dict = require_auth()):
    conn = get_db()
    try:
        where, args = _user_scope_clause(conn, current_user)
        items = []
        for row in conn.execute(
            f"SELECT * FROM translations WHERE 1=1{where} ORDER BY created_at DESC LIMIT 50", args
        ).fetchall():
            items.append(dict(row))
        return items
    finally:
        conn.close()


@router.delete("/api/translation/{task_id}")
async def delete_translation(task_id: str, current_user: dict = require_auth()):
    conn = get_db()
    try:
        where, args = _user_scope_clause(conn, current_user)
        conn.execute(f"DELETE FROM translations WHERE id=?{where}", [task_id] + args)
        conn.commit()
        return {"ok": True}
    finally:
        conn.close()


# ── 翻译术语表（v15）：用户自定义术语 → 翻译时强制应用 ──
class GlossaryItemRequest(BaseModel):
    source_term: str = Field(..., min_length=1, max_length=100)
    target_term: str = Field(..., min_length=1, max_length=200)


@router.get("/api/translation/glossary")
async def list_glossary(current_user: dict = require_auth()):
    """用户术语表列表（个人数据，按 user_id 隔离）。"""
    return _load_glossary(str(current_user.get("user_id", "")))


@router.post("/api/translation/glossary")
async def add_glossary_item(data: GlossaryItemRequest, current_user: dict = require_auth()):
    """新增术语条目（同用户同原文已存在时更新译文）。"""
    source = data.source_term.strip()
    target = data.target_term.strip()
    if not source or not target:
        raise HTTPException(400, "术语原文与译文均不能为空")
    user_id = str(current_user.get("user_id", ""))
    conn = get_db()
    try:
        gid = f"glossary_{uuid.uuid4().hex[:10]}"
        conn.execute(
            "INSERT INTO translation_glossary (id, user_id, source_term, target_term) VALUES (?,?,?,?)",
            (gid, user_id, source, target),
        )
        conn.commit()
        return {"ok": True, "id": gid}
    finally:
        conn.close()


@router.delete("/api/translation/glossary/{gid}")
async def delete_glossary_item(gid: str, current_user: dict = require_auth()):
    """删除术语条目（归属校验，仅能删自己的）。"""
    conn = get_db()
    try:
        conn.execute(
            "DELETE FROM translation_glossary WHERE id=? AND user_id=?",
            (gid, str(current_user.get("user_id", ""))),
        )
        conn.commit()
        return {"ok": True}
    finally:
        conn.close()


# ── 双语对照导出（v15）：md / docx ──
TRANSLATION_EXPORT_DIR = os.path.join(os.path.dirname(__file__), "uploads", "translations")
os.makedirs(TRANSLATION_EXPORT_DIR, exist_ok=True)


class TranslationExportRequest(BaseModel):
    source: str = Field(..., max_length=50000)
    translation: str = Field(..., max_length=50000)
    format: str = "md"  # md | docx


def _align_paragraphs(source: str, translation: str) -> list:
    """段落级对齐：行数一致逐行对照，否则整块对照。"""
    src_lines = [l for l in (source or "").splitlines() if l.strip()]
    tgt_lines = [l for l in (translation or "").splitlines() if l.strip()]
    if len(src_lines) == len(tgt_lines) and len(src_lines) > 0:
        return [{"source": s, "translation": t} for s, t in zip(src_lines, tgt_lines)]
    return [{"source": source or "", "translation": translation or ""}]


def _build_bilingual_md(source: str, translation: str) -> str:
    """双语对照 Markdown（段落级原文/译文成对）。"""
    pairs = _align_paragraphs(source, translation)
    lines = ["# 双语对照翻译", ""]
    for i, p in enumerate(pairs, 1):
        lines += [f"## 第 {i} 段", "", "**原文**", "", p["source"], "", "**译文**", "", p["translation"], ""]
    return "\n".join(lines)


def _build_bilingual_docx(source: str, translation: str, path: str) -> None:
    """双语对照 Word 文档。"""
    from docx import Document

    doc = Document()
    doc.add_heading("双语对照翻译", level=0)
    for i, p in enumerate(_align_paragraphs(source, translation), 1):
        doc.add_heading(f"第 {i} 段", level=1)
        doc.add_paragraph(p["source"])
        doc.add_paragraph(p["translation"])
    doc.save(path)


@router.post("/api/translation/export")
async def export_translation(data: TranslationExportRequest, current_user: dict = require_auth()):
    """双语对照导出（md/docx），返回下载地址。"""
    fmt = data.format.lower()
    if fmt not in ("md", "docx"):
        raise HTTPException(400, "仅支持 md/docx 格式")
    fname = f"translation_{uuid.uuid4().hex[:10]}.{fmt}"
    fpath = os.path.join(TRANSLATION_EXPORT_DIR, fname)
    try:
        if fmt == "docx":
            _build_bilingual_docx(data.source, data.translation, fpath)
        else:
            with open(fpath, "w", encoding="utf-8") as wf:
                wf.write(_build_bilingual_md(data.source, data.translation))
        return {"ok": True, "filename": fname, "download_url": f"/api/translation/download/{fname}"}
    except Exception as e:
        logger.exception("translation export failed")
        raise HTTPException(500, "操作失败，请稍后重试") from e


@router.get("/api/translation/download/{filename}")
async def download_translation_export(filename: str, current_user: dict = require_auth()):
    """下载导出的双语对照文件（basename 白名单防目录穿越）。"""
    safe = os.path.basename(filename)
    path = os.path.join(TRANSLATION_EXPORT_DIR, safe)
    if not os.path.exists(path):
        raise HTTPException(404, "文件不存在或已过期清理")
    media = (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        if safe.endswith(".docx")
        else "text/markdown"
    )
    return FileResponse(path, filename=safe, media_type=media)


# ══════════════════════════════════════════════════════════════
# Phase 4: 运营分析
# ══════════════════════════════════════════════════════════════


class ABTestCreate(BaseModel):
    name: str
    description: str = ""
    variant_a: str = ""
    variant_b: str = ""


@router.get("/api/dashboard/stats")
async def get_dashboard_stats(current_user: dict = require_auth()):
    """获取仪表盘统计数据"""
    conn = get_db()
    try:
        stats = {
            "agents": conn.execute("SELECT COUNT(*) FROM agents WHERE active=1").fetchone()[0],
            "workflows": conn.execute("SELECT COUNT(*) FROM workflows WHERE active=1").fetchone()[0],
            "projects": conn.execute("SELECT COUNT(*) FROM projects WHERE active=1").fetchone()[0],
            "tasks": conn.execute("SELECT COUNT(*) FROM global_tasks WHERE active=1").fetchone()[0],
            "tasks_completed": conn.execute(
                "SELECT COUNT(*) FROM global_tasks WHERE status='done' AND active=1"
            ).fetchone()[0],
            "pipelines": conn.execute("SELECT COUNT(*) FROM pipelines WHERE active=1").fetchone()[0],
            "code_generations": conn.execute("SELECT COUNT(*) FROM code_generations").fetchone()[0],
            "translations": conn.execute("SELECT COUNT(*) FROM translations").fetchone()[0],
            "artifacts": conn.execute("SELECT COUNT(*) FROM artifacts WHERE active=1").fetchone()[0],
        }
        return stats
    finally:
        conn.close()


@router.get("/api/analytics/overview")
async def get_analytics_overview(current_user: dict = require_auth()):
    """获取分析概览"""
    conn = get_db()
    try:
        return {
            "total_agents": conn.execute("SELECT COUNT(*) FROM agents WHERE active=1").fetchone()[0],
            "total_workflows": conn.execute("SELECT COUNT(*) FROM workflows WHERE active=1").fetchone()[0],
            "total_projects": conn.execute("SELECT COUNT(*) FROM projects WHERE active=1").fetchone()[0],
            "total_tasks": conn.execute("SELECT COUNT(*) FROM global_tasks WHERE active=1").fetchone()[0],
            "completed_tasks": conn.execute("SELECT COUNT(*) FROM global_tasks WHERE status='done'").fetchone()[0],
            "total_artifacts": conn.execute("SELECT COUNT(*) FROM artifacts WHERE active=1").fetchone()[0],
            "total_code_gens": conn.execute("SELECT COUNT(*) FROM code_generations").fetchone()[0],
            "total_translations": conn.execute("SELECT COUNT(*) FROM translations").fetchone()[0],
        }
    finally:
        conn.close()


@router.get("/api/ab-tests")
async def list_ab_tests(current_user: dict = require_auth()):
    conn = get_db()
    try:
        items = []
        for row in conn.execute("SELECT * FROM ab_tests WHERE active=1 ORDER BY created_at DESC").fetchall():
            t = dict(row)
            t["result"] = json.loads(t.get("result", "{}"))
            items.append(t)
        return items
    finally:
        conn.close()


@router.post("/api/ab-tests")
async def create_ab_test(data: ABTestCreate, current_user: dict = require_auth()):
    conn = get_db()
    try:
        tid = f"ab_{uuid.uuid4().hex[:12]}"
        conn.execute(
            "INSERT INTO ab_tests (id, name, description, variant_a, variant_b, created_by) VALUES (?,?,?,?,?,?)",
            (tid, data.name, data.description, data.variant_a, data.variant_b, current_user["username"]),
        )
        conn.commit()
        return {"ok": True, "id": tid}
    finally:
        conn.close()


@router.delete("/api/ab-tests/{tid}")
async def delete_ab_test(tid: str, current_user: dict = require_auth()):
    conn = get_db()
    try:
        conn.execute("UPDATE ab_tests SET active=0 WHERE id=?", (tid,))
        conn.commit()
        return {"ok": True}
    finally:
        conn.close()


class ABTestRunRequest(BaseModel):
    objective: str = "整体效果"  # 实验目标维度，如：标题吸引力 / 转化率 / 用户偏好


# AB 实验评分维度（营销专家口径，与结果展示前端一一对应）
_AB_DIMENSIONS = ["吸引力", "清晰度", "转化力", "专业度", "记忆点"]

_AB_RUN_SYSTEM = """你是拥有 20 年经验的增长实验与文案转化专家，擅长 A/B 测试设计、
方案评估与数据化决策。请基于给定的 A/B 两个方案，完成四件事：
1. 将 A、B 方案各自扩写为完整、可落地、有细节的终稿内容（保留原意，增强说服力）；
2. 从五个维度分别对 A、B 打分（0-100 整数）：吸引力、清晰度、转化力、专业度、记忆点；
3. 给出明确胜出方、置信度（0-100）与一句话结论建议；
4. 给出结构化分析：胜出原因（结合维度分差说明为什么赢）、风险提示（采用胜出方案的潜在风险）与下一步行动（验证/发布建议）。
只返回 JSON，不要任何其他文字。"""

_AB_RUN_USER = """实验名称：{name}
实验目标：{objective}
方案 A：{variant_a}
方案 B：{variant_b}

请严格按此 JSON 结构返回：
{{
  "generated_a": "方案A扩写后的完整终稿",
  "generated_b": "方案B扩写后的完整终稿",
  "scores": [{{"dimension": "吸引力", "a": 0, "b": 0}}],
  "winner": "A" | "B",
  "confidence": 0,
  "conclusion": "一句话结论与建议",
  "analysis": {{
    "winner_reason": "胜出原因分析（结合维度分差）",
    "risks": ["风险1", "风险2"],
    "next_steps": ["下一步行动1", "下一步行动2"]
  }}
}}"""





def _normalize_ab_scores(parsed: dict) -> list:
    """规范化 AB 五维评分：钳制 0-100、补缺失维度、保留额外维度。"""
    scores_map: dict = {}
    for s in parsed.get("scores") or []:
        if not isinstance(s, dict) or not s.get("dimension"):
            continue
        dim = s["dimension"]
        scores_map[dim] = {"dimension": dim, "a": _ab_num(s.get("a")), "b": _ab_num(s.get("b"))}
    scores = [scores_map.get(d, {"dimension": d, "a": 0, "b": 0}) for d in _AB_DIMENSIONS]
    for dim, s in scores_map.items():
        if dim not in _AB_DIMENSIONS:
            scores.append(s)
    return scores


def _ab_num(v) -> int:
    """AB 维度分数钳制 0-100。"""
    try:
        return max(0, min(100, int(v)))
    except (ValueError, TypeError):
        return 0


def _ab_infer_winner(scores: list, winner: str | None) -> str:
    """按总分推断胜出方。"""
    if winner in ("A", "B"):
        return winner
    total_a = sum(s["a"] for s in scores)
    total_b = sum(s["b"] for s in scores)
    return "A" if total_a >= total_b else "B"


def _ab_winner_reason(scores: list, winner: str, total_a: int, total_b: int, reason: str) -> str:
    """胜出原因：优先 LLM 输出，否则用分差最大维度派生。"""
    if reason:
        return reason
    diffs = sorted(scores, key=lambda s: abs(s["a"] - s["b"]), reverse=True)
    if diffs and diffs[0]["a"] != diffs[0]["b"]:
        top = diffs[0]
        win_side = top["a"] if winner == "A" else top["b"]
        lose_side = top["b"] if winner == "A" else top["a"]
        return (
            f"方案 {winner} 在「{top['dimension']}」维度领先 {win_side - lose_side} 分，"
            f"五维总分亦占优（{total_a} vs {total_b}），故判定胜出。"
        )
    return f"方案 {winner} 五维评分更高，判定胜出。"

def normalize_ab_result(parsed: dict, objective: str = "") -> dict:
    """AB 实验结果结构化兜底（纯函数，可单测）。

    LLM 输出可能缺失维度分/胜出方/置信度/分析段，此处补齐五维空分、
    按总分推断胜出方、按分差推断置信度，并派生胜出原因，保证前端渲染字段齐全。
    """
    parsed = parsed or {}

    # 维度分：非法值钳制 0-100；缺失维度补 0 分；LLM 额外维度保留
    scores = _normalize_ab_scores(parsed)
    total_a = sum(s["a"] for s in scores)
    total_b = sum(s["b"] for s in scores)
    winner = _ab_infer_winner(scores, parsed.get("winner"))

    try:
        confidence = max(0, min(100, int(parsed.get("confidence") or 0)))
    except (ValueError, TypeError):
        confidence = 0
    if confidence == 0 and total_a != total_b:
        diff = abs(total_a - total_b)
        confidence = round(min(90, max(50, diff / max(total_a, total_b, 1) * 100)))

    analysis = parsed.get("analysis") if isinstance(parsed.get("analysis"), dict) else {}
    winner_reason = (analysis.get("winner_reason") or "").strip()
    risks = [str(r).strip() for r in (analysis.get("risks") or []) if str(r).strip()]
    next_steps = [str(n).strip() for n in (analysis.get("next_steps") or []) if str(n).strip()]
    winner_reason = _ab_winner_reason(scores, winner, total_a, total_b, winner_reason)

    return {
        "status": "completed",
        "objective": (objective or "整体效果").strip()[:40],
        "generated_a": parsed.get("generated_a", ""),
        "generated_b": parsed.get("generated_b", ""),
        "scores": scores,
        "winner": winner,
        "confidence": confidence,
        "conclusion": (parsed.get("conclusion") or "").strip(),
        "analysis": {
            "winner_reason": winner_reason,
            "risks": risks,
            "next_steps": next_steps,
        },
        "ran_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    }




def _ab_md_sections(test: dict, result: dict) -> list:
    """A/B 报告各段落行（背景/方案/评分/结论/分析）。"""
    lines = []
    if test.get("description"):
        lines += ["## 实验背景", "", test["description"], ""]
    lines += ["## 方案对比", "", f"**方案 A**：{test.get('variant_a') or '-'}", "", f"**方案 B**：{test.get('variant_b') or '-'}", ""]
    if result.get("generated_a") or result.get("generated_b"):
        lines.append("## AI 扩写终稿")
        if result.get("generated_a"):
            lines += ["", "### 方案 A", "", result["generated_a"], ""]
        if result.get("generated_b"):
            lines += ["", "### 方案 B", "", result["generated_b"], ""]
    scores = result.get("scores") or []
    if scores:
        lines += ["## 五维评分对比", "", "| 维度 | 方案A | 方案B |", "| --- | ---: | ---: |"]
        lines += [f"| {s.get('dimension') or '-'} | {s.get('a') or 0} | {s.get('b') or 0} |" for s in scores]
        lines.append("")
    return lines


def _ab_md_analysis(result: dict) -> list:
    """A/B 报告分析段（胜出原因/风险/下一步）。"""
    lines = []
    analysis = result.get("analysis") if isinstance(result.get("analysis"), dict) else {}
    if analysis.get("winner_reason"):
        lines += ["## 胜出原因", "", analysis["winner_reason"], ""]
    if analysis.get("risks"):
        lines += ["## 风险提示", ""]
        lines += [f"- {r}" for r in analysis["risks"]]
        lines.append("")
    if analysis.get("next_steps"):
        lines += ["## 下一步行动", ""]
        lines += [f"- {n}" for n in analysis["next_steps"]]
        lines.append("")
    return lines

def build_ab_report_md(test: dict, result: dict) -> str:
    """A/B 实验报告 → Markdown（纯函数，可单测；用于报告导出）。"""
    result = result or {}
    lines = [
        "# A/B 测试分析报告",
        "",
        f"- 实验名称：{test.get('name') or '-'}",
        f"- 实验目标：{result.get('objective') or '-'}",
        f"- 运行时间：{result.get('ran_at') or '-'}",
        "",
    ]
    lines += _ab_md_sections(test, result)
    lines += [
        "## 结论",
        "",
        f"- 胜出方：方案 {result.get('winner') or '-'}",
        f"- 置信度：{result.get('confidence') or 0}%",
        f"- 决策建议：{result.get('conclusion') or '-'}",
        "",
    ]
    lines += _ab_md_analysis(result)
    lines += ["---", "由AI 星火 AI A/B 测试生成"]
    return "\n".join(lines)


@router.post("/api/ab-tests/{tid}/run")
async def run_ab_test(tid: str, data: ABTestRunRequest, current_user: dict = require_auth()):
    """运行 A/B 实验：LLM 分别扩写 A/B 方案并按 5 维度打分，产出胜出方与置信度。

    结果落库 ab_tests.result（JSON），供结果页直接读取；同步返回完整结果。
    """
    conn = get_db()
    try:
        row = conn.execute("SELECT * FROM ab_tests WHERE id=? AND active=1", (tid,)).fetchone()
    finally:
        conn.close()
    if not row:
        raise HTTPException(404, "实验不存在")
    if not (row["variant_a"] or row["variant_b"]):
        raise HTTPException(400, "A/B 方案为空，无法运行")

    user_prompt = _AB_RUN_USER.format(
        name=row["name"],
        objective=(data.objective or "整体效果").strip()[:40],
        variant_a=row["variant_a"] or "（未设置）",
        variant_b=row["variant_b"] or "（未设置）",
    )
    payload = await call_llm_async(_AB_RUN_SYSTEM, user_prompt, json_mode=True)
    parsed = parse_llm_json(payload)

    result = normalize_ab_result(parsed, objective=data.objective)
    conn = get_db()
    try:
        conn.execute(
            "UPDATE ab_tests SET status='completed', result=?, updated_at=? WHERE id=?",
            (json.dumps(result, ensure_ascii=False), datetime.now().isoformat(), tid),
        )
        conn.commit()
    finally:
        conn.close()
    return result


@router.get("/api/ab-tests/{tid}/results")
async def get_ab_test_results(tid: str, current_user: dict = require_auth()):
    """读取实验运行结果（未运行时返回 status=pending）。"""
    conn = get_db()
    try:
        row = conn.execute("SELECT result, status FROM ab_tests WHERE id=? AND active=1", (tid,)).fetchone()
    finally:
        conn.close()
    if not row:
        raise HTTPException(404, "实验不存在")
    try:
        result = json.loads(row["result"] or "{}")
    except Exception:
        result = {}
    if not result.get("status"):
        result = {"status": "pending"}
    return result


@router.get("/api/ab-tests/{tid}/report")
async def get_ab_test_report(tid: str, current_user: dict = require_auth()):
    """导出 A/B 实验分析报告（Markdown）；实验不存在或未运行返回 404。"""
    conn = get_db()
    try:
        row = conn.execute("SELECT * FROM ab_tests WHERE id=? AND active=1", (tid,)).fetchone()
    finally:
        conn.close()
    if not row:
        raise HTTPException(404, "实验不存在")
    try:
        result = json.loads(row["result"] or "{}")
    except Exception:
        result = {}
    if not result.get("status"):
        raise HTTPException(404, "实验尚未运行")
    test = {
        "name": row["name"],
        "description": row["description"],
        "variant_a": row["variant_a"],
        "variant_b": row["variant_b"],
    }
    content = build_ab_report_md(test, result)
    base = row["name"] or "ab-test"
    return {"filename": f"{base}-AB实验报告.md", "content": content}


# ══════════════════════════════════════════════════════════════
# 办公效率: PPT + Excel
# ══════════════════════════════════════════════════════════════


class PPTGenerateRequest(BaseModel):
    title: str
    outline: str = ""
    model: str = ""
    template: str = "business"


class ExcelRequest(BaseModel):
    operation: str = "create"
    title: str = ""
    data: dict = {}


# ── PPT 模板库（v15）：4 类场景模板 = 色板 + 结构原则 + 字体 ──
PPT_TEMPLATES = {
    "business": {
        "name": "商务汇报",
        "desc": "面向管理层的工作汇报/项目提案",
        "palette": {
            "dark": (0x1B, 0x26, 0x3B),
            "accent": (0x4F, 0x46, 0xE5),
            "accent_light": (0xEE, 0xF0, 0xFF),
            "gray": (0x6B, 0x72, 0x80),
            "text": (0x33, 0x3A, 0x4A),
            "white": (0xFF, 0xFF, 0xFF),
        },
        "font": "Microsoft YaHei",
        "principles": (
            "1. 黄金结构：封面→目录→背景/问题（Why）→方案/举措（What）→执行/数据（How）→案例/成果→下一步→感谢页\n"
            "2. 每页一个结论式标题（一句话结论，非描述性标题）\n"
            "3. 数据说话：观点用可量化数据支撑，标注数据来源\n"
            "4. 视觉思维：对比用柱状图、趋势用折线图、占比用饼图\n"
            "5. 语言风格：严谨克制、决策导向，避免空话套话"
        ),
    },
    "roadshow": {
        "name": "融资路演",
        "desc": "面向投资人的融资路演/产品发布",
        "palette": {
            "dark": (0x16, 0x0E, 0x2B),
            "accent": (0xE1, 0x1D, 0x48),
            "accent_light": (0xFD, 0xEC, 0xF0),
            "gray": (0x8A, 0x84, 0x9A),
            "text": (0x2A, 0x24, 0x3B),
            "white": (0xFF, 0xFF, 0xFF),
        },
        "font": "Microsoft YaHei",
        "principles": (
            "1. 故事线：痛点→解决方案→市场机会（TAM/SAM/SOM）→产品→商业模式→竞争壁垒→团队→里程碑→融资需求\n"
            "2. 市场机会量化：市场规模、增长率、可服务市场占比，引用权威数据源\n"
            "3. 商业模式清晰：收入模型、单位经济模型（LTV/CAC）、毛利结构\n"
            "4. 竞争格局：对比表展示差异化壁垒（技术/渠道/数据/成本）\n"
            "5. 语言风格：激情与克制并重，每页回答投资人最关心的一个疑问"
        ),
    },
    "teaching": {
        "name": "教学课件",
        "desc": "面向学员的知识培训/课堂课件",
        "palette": {
            "dark": (0x0F, 0x33, 0x24),
            "accent": (0x0E, 0x9F, 0x6E),
            "accent_light": (0xE6, 0xF7, 0xF1),
            "gray": (0x6B, 0x72, 0x80),
            "text": (0x1F, 0x29, 0x37),
            "white": (0xFF, 0xFF, 0xFF),
        },
        "font": "Microsoft YaHei",
        "principles": (
            "1. 知识结构：导入（为什么学）→概念定义→原理讲解→案例演示→练习/互动→小结复习\n"
            "2. 每页聚焦一个知识点，概念先行、实例紧随\n"
            "3. 善用类比与图示降低理解门槛，关键术语加粗并给出通俗解释\n"
            "4. 设计互动环节：提问、随堂练习、小组讨论建议\n"
            "5. 语言风格：亲切易懂、循序渐进，避免堆砌术语"
        ),
    },
    "marketing": {
        "name": "营销方案",
        "desc": "面向团队的营销策划/活动方案",
        "palette": {
            "dark": (0x2B, 0x16, 0x3B),
            "accent": (0xDB, 0x27, 0x77),
            "accent_light": (0xFD, 0xEF, 0xF7),
            "gray": (0x7A, 0x70, 0x80),
            "text": (0x35, 0x2B, 0x3D),
            "white": (0xFF, 0xFF, 0xFF),
        },
        "font": "Microsoft YaHei",
        "principles": (
            "1. 方案主线：市场洞察（用户/竞争/趋势）→策略核心（定位/主张/差异化 ）→创意内容（话题/物料/渠道）→执行节奏（甘特/里程碑）→预算与KPI→风险预案\n"
            "2. 洞察先行：每个策略决策附用户痛点/数据证据\n"
            "3. 创意具体化：给出可直接执行的文案方向、视觉风格、投放组合\n"
            "4. 结果导向：设定可量化的 KPI（曝光/转化/ROI）与达成路径\n"
            "5. 语言风格：生动有画面感，策略严谨、创意大胆"
        ),
    },
    "tech": {
        "name": "科技产品",
        "desc": "技术方案/产品发布/开发者分享",
        "palette": {
            "dark": (0x0A, 0x18, 0x2E),
            "accent": (0x00, 0x96, 0xF7),
            "accent_light": (0xE8, 0xF5, 0xFE),
            "gray": (0x64, 0x74, 0x8B),
            "text": (0x1E, 0x29, 0x3B),
            "white": (0xFF, 0xFF, 0xFF),
        },
        "font": "Microsoft YaHei",
        "principles": (
            "1. 技术叙事：背景/痛点→技术方案→架构设计→关键技术亮点→性能与效果→落地路径\n"
            "2. 架构图优先：用分层/流程/拓扑描述系统，避免大段代码\n"
            "3. 指标量化：性能提升、成本下降、延迟降低用具体数字对比（Before/After）\n"
            "4. 前瞻性：给出技术趋势判断与演进路线图\n"
            "5. 语言风格：理性克制、术语准确，为开发者与决策者双向服务"
        ),
    },
    "consulting": {
        "name": "咨询分析",
        "desc": "行业研究/战略咨询/尽调报告",
        "palette": {
            "dark": (0x12, 0x1A, 0x22),
            "accent": (0xC2, 0x41, 0x0C),
            "accent_light": (0xFD, 0xF0, 0xE7),
            "gray": (0x6B, 0x72, 0x80),
            "text": (0x2D, 0x37, 0x42),
            "white": (0xFF, 0xFF, 0xFF),
        },
        "font": "Microsoft YaHei",
        "principles": (
            "1. 金字塔结构：结论先行→关键论据（2-4个）→支撑细节/数据→洞察与建议\n"
            "2. MECE 框架：每一层的论点互斥且穷尽，逻辑无遗漏\n"
            "3. 数据驱动：市场规模/增速/份额等关键数字标注来源与测算口径\n"
            "4. 可执行建议：每个洞察附落地动作与预期影响（What-Why-How）\n"
            "5. 语言风格：客观中立、分析严谨，避免主观断言"
        ),
    },
    "finance": {
        "name": "金融投研",
        "desc": "投资分析/行业研究/资金汇报",
        "palette": {
            "dark": (0x14, 0x1E, 0x2E),
            "accent": (0x0E, 0x93, 0x84),
            "accent_light": (0xE6, 0xF6, 0xF4),
            "gray": (0x64, 0x70, 0x84),
            "text": (0x22, 0x2B, 0x3D),
            "white": (0xFF, 0xFF, 0xFF),
        },
        "font": "Microsoft YaHei",
        "principles": (
            "1. 投资逻辑链：行业空间→竞争格局→商业模式→财务表现→估值与风险→投资建议\n"
            "2. 财务三表核心指标：营收/利润/现金流/毛利/净利率，同比环比双维度\n"
            "3. 风险揭示：政策/市场/技术/经营风险逐一列出，附应对措施\n"
            "4. 估值方法：PE/PEG/DCF 多方法交叉验证，给出假设与敏感性\n"
            "5. 语言风格：审慎专业、结论明确，符合合规要求（不承诺收益）"
        ),
    },
}


# ── PPT 商业化：专家级大纲 prompt（异步任务 worker 复用）──
def _build_ppt_system_prompt(template_id: str = "business") -> str:
    """按模板生成 PPT 大纲 System Prompt（模板结构原则 + 段落级结构化输出要求）。"""
    tpl = PPT_TEMPLATES.get(template_id) or PPT_TEMPLATES["business"]
    return f"""你是资深PPT策划与演示设计专家，服务于500强企业高管汇报场景。

## 模板：{tpl['name']}（{tpl['desc']}）

### 结构设计原则
{tpl['principles']}

## 段落级结构化要求
每页 content 使用段落级结构（level 分层 + emphasis 强调），便于渲染为层级清晰的版面：
- level 0：主论点（结论式短句，一页仅1条）
- level 1：支撑论据/细节（每页3-5条，每条必须含具体数据、场景或案例，禁止空泛表述）
- emphasis：strong（关键数字或结论，渲染加粗高亮）、quote（引用/金句，渲染为引用样式）、normal（普通）

## 内容丰富度要求（每页必须满足）
1. 每页 content 至少 4 条：1 条 level 0 结论 + 3-5 条 level 1 支撑；每条论据必须落到具体数据（百分比/金额/倍数/年份）、可感知场景或真实案例，禁止"提升效率""加强协同"类空话
2. 数据页（type=data）：chart_suggestion 必须写明具体数据点（如"近三年营收：1.2亿→2.3亿→4.1亿，柱状图对比"），预览器据此渲染真实维度图表
3. 案例页（type=case）：content 必须按"背景 → 行动 → 结果"三段组织（3 条 level 1），关键数字用 emphasis=strong 标注
4. 内容页（type=content）：遵循"结论 + 论据×3左右 + 关键数据 + 行动建议"结构，关键数据与结论用 emphasis=strong 标注
5. 所有页面必须填写 subtitle（副标题/上下文说明）、notes（演讲备注：过渡语+强调点+互动问题）、duration_seconds（预估秒数）
6. 页面总量 10-14 页：cover 1 + toc 1 + content 4-7 + data 2-3 + case 1-2 + summary 1 + thanks 1

## 通用设计原则
1. 每页原则：一页一个核心观点，标题即是结论（非描述性标题）
2. 视觉思维：优先用图表代替文字（对比用柱状图、趋势用折线图、占比用饼图、流程用箭头）
3. 记忆锚点：设计一个贯穿全篇的视觉隐喻或故事线

请严格按以下JSON格式返回（只返回JSON，不要任何其他文字）：
{{
  "meta": {{
    "storyline": "一句话概括全篇叙述逻辑",
    "visual_theme": "建议配色方案和视觉风格",
    "estimated_duration": "预计演讲时长（分钟）"
  }},
  "slides": [
    {{
      "type": "cover|toc|content|data|case|summary|thanks",
      "title": "结论式标题（不超过15字）",
      "subtitle": "副标题或上下文说明",
      "content": [
        {{"text": "主论点", "level": 0, "emphasis": "strong"}},
        {{"text": "支撑论据", "level": 1, "emphasis": "normal"}}
      ],
      "chart_suggestion": "如适用，建议的图表类型和数据维度",
      "notes": "演讲备注：过渡语、强调点、互动问题",
      "duration_seconds": 60
    }}
  ]
}}

生成10-14页幻灯片，确保逻辑递进、首尾呼应、信息密度饱满（每页都能独立成稿）。"""


def _parse_ppt_outline(result: str) -> dict:
    """LLM 大纲文本 → dict（容错：去代码块包裹/截取首个 JSON/兜底空结构）"""
    text = (result or "").strip()
    text = re.sub(r"^```(?:json)?\s*|\s*```$", "", text)
    start, end = text.find("{"), text.rfind("}")
    if start >= 0 and end > start:
        text = text[start : end + 1]
    try:
        data = parse_llm_json(text)
    except Exception:
        return {"meta": {}, "slides": []}
    if not isinstance(data, dict):
        return {"meta": {}, "slides": []}
    data.setdefault("meta", {})
    if not isinstance(data.get("slides"), list):
        data["slides"] = []
    return data



def _add_pptx_slide(slide_layout, title: str, content: list) -> None:
    """添加幻灯片内容。"""
    from pptx.util import Inches
    title_box = slide_layout.shapes.title
    content_box = slide_layout.shapes.placeholders[1]
    if title_box:
        title_box.text = title
    if content_box:
        content_box.text = "\n".join(content)

def _format_pptx_content(outline: dict) -> list:
    """格式化PPTX内容。"""
    sections = []
    for section in outline.get("sections", []):
        sections.append({
            "title": section.get("title", ""),
            "content": section.get("content", [])
        })
    return sections


def _prepare_pptx_context(presentation_data):
    """准备PPT构建上下文。"""
    return {
        "data": presentation_data,
        "slides": [],
        "status": "prepared"
    }

def _build_single_slide(slide_index, slide_data):
    """构建单个幻灯片。"""
    return {
        "index": slide_index,
        "data": slide_data,
        "status": "built"
    }

def _finalize_pptx_result(slides):
    """汇总PPT构建结果。"""
    return {
        "total_slides": len(slides),
        "slides": slides,
        "status": "completed"
    }


def _build_pptx_simple_v2(slides_data: list, output_path: str) -> str:
    """简化版PPT构建。"""
    try:
        from pptx import Presentation
        prs = Presentation()
        
        for slide_data in slides_data:
            prs.slides.add_slide(prs.slide_layouts[6])
            title = slide_data.get("title", "Slide")
            prs.slides[-1].shapes.title.text = title
        
        prs.save(output_path)
        return output_path
    except Exception as e:
        logger.info(f"PPT构建失败: {e}")
        return ""


def _pptx_cover(slide, title_text: str, subtitle: str, meta: dict, rect, text, colors) -> None:
    """封面版式渲染。"""
    from pptx.enum.text import PP_ALIGN

    rect(slide, 0, 0, 13.333, 7.5, colors["DARK"])
    rect(slide, 0, 4.7, 13.333, 0.06, colors["ACCENT"])
    text(slide, 1.2, 2.6, 10.9, 1.2, title_text, 40, colors["WHITE"], bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        text(slide, 1.2, 3.9, 10.9, 0.6, subtitle, 20, colors["GRAY"], align=PP_ALIGN.CENTER, first=False)
    text(slide, 1.2, 6.8, 10.9, 0.4,
         f"预计时长 {meta.get('estimated_duration', '-')} 分钟  |  视觉主题：{meta.get('visual_theme', '-')}",
         12, colors["GRAY"], align=PP_ALIGN.CENTER)


def _pptx_thanks(slide, title_text: str, subtitle: str, rect, text, colors) -> None:
    """致谢版式渲染。"""
    from pptx.enum.text import PP_ALIGN

    rect(slide, 0, 0, 13.333, 7.5, colors["DARK"])
    text(slide, 1.2, 3.1, 10.9, 1.0, title_text, 44, colors["WHITE"], bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        text(slide, 1.2, 4.2, 10.9, 0.6, subtitle, 20, colors["GRAY"], align=PP_ALIGN.CENTER)


def _pptx_toc(slide, title_text: str, content: list, rect, text, colors) -> None:
    """目录版式渲染。"""
    rect(slide, 0, 0, 13.333, 7.5, colors["WHITE"])
    rect(slide, 0, 0, 0.25, 7.5, colors["ACCENT"])
    text(slide, 0.8, 0.6, 8, 0.8, "目录", 30, colors["DARK"], bold=True)
    for j, item in enumerate(content):
        text(slide, 1.0, 1.9 + j * 0.95, 10, 0.7, f"{j + 1:02d}    {item}", 20, colors["TEXT"])


def _pptx_case(slide, title_text: str, content: list, chart_suggestion: str, rect, text, bullets, colors) -> None:
    """案例版式渲染。"""
    rect(slide, 0, 0, 13.333, 7.5, colors["ACCENT_LIGHT"])
    rect(slide, 0, 0, 13.333, 0.18, colors["ACCENT"])
    text(slide, 0.8, 0.55, 11.7, 0.7, title_text, 26, colors["DARK"], bold=True)
    bullets(slide, content, 0.8, 1.6, 11.7, 5.0, size=18, gap=12)
    if chart_suggestion:
        text(slide, 0.8, 6.4, 11.7, 0.5, f"📊 可视化建议：{chart_suggestion}", 13, colors["GRAY"])


def _pptx_data(slide, title_text: str, subtitle: str, content: list, chart_suggestion: str, rect, text, bullets, colors) -> None:
    """数据版式渲染。"""
    rect(slide, 0, 0, 13.333, 7.5, colors["WHITE"])
    rect(slide, 0, 0, 13.333, 0.18, colors["ACCENT"])
    text(slide, 0.8, 0.55, 11.7, 0.7, title_text, 26, colors["DARK"], bold=True)
    if subtitle:
        text(slide, 0.8, 1.25, 11.7, 0.45, subtitle, 14, colors["GRAY"])
    bullets(slide, content, 0.8, 1.95, 11.7, 4.2, size=18, gap=12)
    if chart_suggestion:
        rect(slide, 0.8, 6.15, 11.7, 0.75, colors["ACCENT_LIGHT"])
        text(slide, 1.05, 6.3, 11.2, 0.45, f"📊 数据可视化：{chart_suggestion}", 14, colors["ACCENT"], bold=True)


def _pptx_content(slide, title_text: str, subtitle: str, content: list, chart_suggestion: str, rect, text, bullets, colors) -> None:
    """内容版式渲染（默认）。"""
    rect(slide, 0, 0, 13.333, 7.5, colors["WHITE"])
    rect(slide, 0, 0, 0.25, 7.5, colors["ACCENT"])
    text(slide, 0.8, 0.55, 11.7, 0.7, title_text, 26, colors["DARK"], bold=True)
    if subtitle:
        text(slide, 0.8, 1.25, 11.7, 0.45, subtitle, 14, colors["GRAY"])
    bullets(slide, content, 0.8, 1.95, 11.7, 4.6, size=18, gap=12)
    if chart_suggestion:
        text(slide, 0.8, 6.55, 11.7, 0.45, f"📊 可视化建议：{chart_suggestion}", 13, colors["GRAY"])


def _pptx_make_renderers(tpl: dict):
    """创建 PPT 版式渲染器集合（闭包绑定主题色板/字体）。返回 (helpers, renderers)。"""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    pal = tpl["palette"]
    colors = {
        "DARK": RGBColor(*pal["dark"]),
        "ACCENT": RGBColor(*pal["accent"]),
        "ACCENT_LIGHT": RGBColor(*pal["accent_light"]),
        "GRAY": RGBColor(*pal["gray"]),
        "TEXT": RGBColor(*pal["text"]),
        "WHITE": RGBColor(*pal["white"]),
    }

    def rect(slide, left, top, w, h, color):
        from pptx.enum.shapes import MSO_SHAPE

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def text(slide, left, top, w, h, content, size, color, bold=False, align=PP_ALIGN.LEFT, first=True):
        from pptx.util import Pt

        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = content
        f = run.font
        f.name = tpl["font"]
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
        return box

    def bullets(slide, items: list, left, top, w, h, size=18, color=None, gap=10):
        from pptx.util import Pt

        color = color or colors["TEXT"]
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(gap)
            if isinstance(item, dict):
                content = str(item.get("text", ""))
                level = int(item.get("level", 0) or 0)
                emphasis = str(item.get("emphasis", "normal"))
            else:
                content = str(item)
                level = 0
                emphasis = "normal"
            if not content:
                continue
            indent = "  " if level >= 1 else ""
            prefix = "•  " if level == 0 else "–  "
            run = p.add_run()
            run.text = f"{indent}{prefix}{content}"
            f = run.font
            f.name = tpl["font"]
            f.size = Pt(max(size - (3 if level >= 1 else 0), 10))
            if emphasis == "strong":
                f.bold = True
                f.color.rgb = colors["ACCENT"]
            elif emphasis == "quote":
                f.italic = True
                f.color.rgb = colors["GRAY"]
            else:
                f.bold = False
                f.color.rgb = color
        return box

    def notes(slide, content: str):
        if content:
            slide.notes_slide.notes_text_frame.text = content

    renderers = {
        "cover": lambda s, t, sub, c, cs=None: _pptx_cover(s, t, sub, c, rect, text, colors),
        "thanks": lambda s, t, sub, c, cs=None: _pptx_thanks(s, t, sub, rect, text, colors),
        "toc": lambda s, t, sub, c, cs=None: _pptx_toc(s, t, c, rect, text, colors),
        "case": lambda s, t, sub, c, cs=None: _pptx_case(s, t, c, cs, rect, text, bullets, colors),
        "data": lambda s, t, sub, c, cs=None: _pptx_data(s, t, sub, c, cs, rect, text, bullets, colors),
        "content": lambda s, t, sub, c, cs=None: _pptx_content(s, t, sub, c, cs, rect, text, bullets, colors),
    }

    def page_footer(slide, page_no: int, total: int, dark: bool = False):
        fg = RGBColor(208, 213, 221) if dark else colors["GRAY"]
        text(slide, 0.8, 7.05, 4, 0.35, f"AI 星火 · {tpl['name']}", 10, fg)
        text(slide, 11.6, 7.05, 1.0, 0.35, f"{page_no:02d} / {total:02d}", 10, fg, align=PP_ALIGN.RIGHT)

    return {"rect": rect, "text": text, "bullets": bullets, "notes": notes}, renderers, page_footer



def _build_pptx_file(title: str, outline: dict, template: str = "business") -> str:
    """大纲 dict → 16:9 PPTX 文件（封面/目录/内容/数据/案例/总结/致谢 + 演讲备注），返回保存路径。"""
    from pptx import Presentation
    from pptx.util import Inches

    tpl = PPT_TEMPLATES.get(template) or PPT_TEMPLATES["business"]
    helpers, renderers, page_footer = _pptx_make_renderers(tpl)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    rect, text, notes = helpers["rect"], helpers["text"], helpers["notes"]

    slides = outline.get("slides") or []
    meta = outline.get("meta") or {}
    for i, s in enumerate(slides):
        if not isinstance(s, dict):
            continue
        stype = s.get("type", "content")
        title_text = s.get("title") or f"第 {i + 1} 页"
        subtitle = s.get("subtitle") or ""
        content = s.get("content") or []
        if isinstance(content, str):
            content = [content]
        slide = prs.slides.add_slide(blank)

        render = renderers.get(stype, renderers["content"])
        # cover/thanks 版式第 4 参为 outline meta（dict），其余为 content（list）
        if stype in ("cover", "thanks"):
            render(slide, title_text, subtitle, meta, s.get("chart_suggestion"))
        else:
            render(slide, title_text, subtitle, content, s.get("chart_suggestion"))
        page_footer(slide, i + 1, len(slides), dark=stype in ("cover", "thanks"))
        notes(slide, s.get("notes") or "")

    if not slides:
        slide = prs.slides.add_slide(blank)
        renderers["cover"](slide, title or "未命名演示", "", {"estimated_duration": "-", "visual_theme": "-"})

    path = os.path.join(PPTX_DIR, f"ppt_{uuid.uuid4().hex[:12]}.pptx")
    prs.save(path)
    return path

async def _ppt_worker(payload: dict, progress: Callable | None = None) -> dict:
    """PPT worker：LLM 大纲 → 解析 → 生成 PPTX 文件 → 记录入库（带用户归属）。"""

    def _report(pct: float, stage: str) -> None:
        _notify_progress(progress, pct, stage)

    _report(5, "规划大纲")
    template = payload.get("template", "business")
    prompt = f"主题：{payload.get('title', '')}"
    if payload.get("outline"):
        prompt += f"\n大纲：{payload['outline']}"
    _report(30, "AI 生成大纲中")
    result = await call_llm_async(_build_ppt_system_prompt(template), prompt)
    _report(55, "解析大纲结构")
    outline_data = _parse_ppt_outline(result)
    _report(70, "排版 PPTX 文件")
    pptx_path = _build_pptx_file(payload.get("title", "未命名演示"), outline_data, template)
    filename = os.path.basename(pptx_path)
    _report(90, "保存记录")
    ppt_id = f"ppt_{uuid.uuid4().hex[:12]}"
    with get_db_context() as conn:
        conn.execute(
            "INSERT INTO ppt_generations (id, title, outline, slides, result, model, user_id, file_path) VALUES (?,?,?,?,?,?,?,?)",
            (
                ppt_id,
                payload.get("title", ""),
                payload.get("outline", ""),
                json.dumps(outline_data, ensure_ascii=False),
                result,
                payload.get("model", ""),
                payload.get("user_id", ""),
                f"/api/ppt/download/{filename}",
            ),
        )
    _report(100, "完成")
    return {
        "id": ppt_id,
        "result": result,
        "slides": len(outline_data.get("slides", [])),
        "pptx": f"/api/ppt/download/{filename}",
    }


async def _ppt_handler(task_id: str, payload: dict, update: Callable, ctx: dict) -> dict:
    """异步任务处理器：包装 PPT 生成，回报进度。"""
    return await _ppt_worker(payload, progress=update)


@router.post("/api/ppt/generate")
async def generate_ppt(data: PPTGenerateRequest, current_user: dict = require_auth()):
    """AI PPT 生成（异步任务：大纲 + 真实 PPTX 文件输出 / 进度跟踪 / 自动重试）"""
    if not data.title.strip():
        raise HTTPException(400, "PPT 主题不能为空")
    # 占位符拦截：模板残留 [xxx] 占位符不允许直接生成（避免污染历史记录）
    if re.search(r"\[[^\]]+\]", data.title):
        raise HTTPException(400, "PPT 主题中还有未填写的占位符（如 [产品名]），请先替换为实际内容")
    if len(data.title.strip()) > 200:
        raise HTTPException(400, "PPT 主题过长（上限 200 字）")
    if data.template not in PPT_TEMPLATES:
        raise HTTPException(400, "操作失败，请稍后重试")
    payload = {
        "title": data.title,
        "outline": data.outline,
        "model": data.model,
        "template": data.template,
        "user_id": str(current_user.get("user_id", "")),
        "username": current_user.get("username", ""),
    }
    task = create_task(
        "ppt_generate",
        payload,
        username=current_user.get("username", ""),
        user_id=str(current_user.get("user_id", "")),
        role=current_user.get("role", ""),
    )
    return {"ok": True, "task_id": task["id"], "status": task["status"]}


@router.get("/api/ppt/download/{filename}")
async def download_ppt(filename: str, current_user: dict = require_auth()):
    """下载生成的 PPTX 文件（basename 白名单防目录穿越）"""
    safe = os.path.basename(filename)
    path = os.path.join(PPTX_DIR, safe)
    if not os.path.exists(path):
        raise HTTPException(404, "文件不存在或已过期清理")
    return FileResponse(
        path, filename=safe, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


@router.get("/api/ppt/history")
async def list_ppt_history(current_user: dict = require_auth()):
    conn = get_db()
    try:
        where, args = _user_scope_clause(conn, current_user)
        items = []
        for row in conn.execute(
            f"SELECT * FROM ppt_generations WHERE 1=1{where} ORDER BY created_at DESC LIMIT 50", args
        ).fetchall():
            items.append(dict(row))
        return items
    finally:
        conn.close()


# ── Excel 异常值检测（v15）：IQR 四分位距法，纯函数无 LLM 依赖 ──
def _percentile(sorted_nums: list, p: float) -> float:
    """线性插值分位数（p ∈ [0,1]）。"""
    if not sorted_nums:
        return 0.0
    k = (len(sorted_nums) - 1) * p
    f = int(k)
    c = min(f + 1, len(sorted_nums) - 1)
    return sorted_nums[f] + (sorted_nums[c] - sorted_nums[f]) * (k - f)


def _detect_outliers(text: str, max_cols: int = 20) -> dict:
    """文本表格数据（TSV/CSV/空格分隔）IQR 异常值检测。

    返回结构化结果：每列上下界 + 异常值明细（行号/值/方向）。
    """
    import re

    lines = [l.rstrip() for l in (text or "").strip().splitlines() if l.strip()]
    if len(lines) < 3:
        return {"success": False, "message": "数据至少需要 3 行（1 行表头 + 2 行数据）"}
    first = lines[0]
    if "\t" in first:
        split = lambda s: s.split("\t")  # noqa: E731
    elif "," in first:
        split = lambda s: s.split(",")  # noqa: E731
    else:
        split = lambda s: re.split(r"\s{2,}", s.strip())  # noqa: E731
    headers = [h.strip() for h in split(first)]
    rows = [split(l) for l in lines[1:]]
    ncols = min(len(headers), max_cols)
    columns = []
    for ci in range(ncols):
        values = []
        for ri, row in enumerate(rows):
            if ci >= len(row):
                continue
            cell = row[ci].strip().replace(",", "") if "," in first else row[ci].strip()
            try:
                values.append((ri + 2, float(cell)))
            except ValueError:
                continue
        if len(values) < 4:  # 数值样本过少无统计意义
            continue
        nums = [v for _, v in values]
        sorted_nums = sorted(nums)
        q1 = _percentile(sorted_nums, 0.25)
        q3 = _percentile(sorted_nums, 0.75)
        iqr = q3 - q1
        if iqr == 0:
            continue
        lower, upper = q1 - 1.5 * iqr, q3 + 1.5 * iqr
        outliers = [
            {"row": ri, "value": v, "direction": "偏高" if v > upper else "偏低"}
            for ri, v in values
            if v < lower or v > upper
        ]
        if outliers:
            columns.append(
                {
                    "name": headers[ci],
                    "count": len(outliers),
                    "lower_bound": round(lower, 4),
                    "upper_bound": round(upper, 4),
                    "outliers": outliers,
                }
            )
    return {
        "success": True,
        "method": "IQR 四分位距法（界 = Q1 − 1.5×IQR / Q3 + 1.5×IQR）",
        "total_rows": len(rows),
        "columns": columns,
        "summary": f"共检测 {len(columns)} 个数值列，发现 {sum(c['count'] for c in columns)} 个异常值",
    }


@router.post("/api/excel/operate")
def excel_operate(data: ExcelRequest, current_user: dict = require_auth()):
    """Excel 操作（analyze/formula/outliers/clean/create）"""
    conn = get_db()
    try:
        op_id = f"excel_{uuid.uuid4().hex[:12]}"
        result = ""

        if data.operation == "analyze":
            system_prompt = """你是资深数据分析师，精通Excel数据分析与商业洞察。

分析框架（按此结构输出）：
## 📊 数据概览
- 数据规模（行数、列数）
- 核心指标（均值、中位数、最大值、最小值）

## 🔍 关键发现（3-5条）
每条发现包含：
- 发现描述（用数据说话）
- 业务含义（这对业务意味着什么）
- 行动建议（接下来该做什么）

## 📈 趋势与模式
- 时间趋势（如有日期列）
- 相关性分析（如有多个数值列）
- 异常值识别

## ⚠️ 风险与机会
- 潜在风险点
- 改善机会点

## 💡 建议的可视化方案
- 图表类型 × 数据维度组合建议（如：柱状图 × 按月份销售额）

要求：语言简洁务实，面向业务决策者，避免技术术语堆砌。"""
            result = call_llm(system_prompt, json.dumps(data.data, ensure_ascii=False))
        elif data.operation == "formula":
            system_prompt = """你是Excel高级公式专家，精通VLOOKUP/XLOOKUP/INDEX-MATCH/SUMIFS/数组公式/Power Query等高级功能。

对用户需求，给出可直接粘贴使用的公式，并附参数表逐项说明。

输出严格JSON（只输出JSON，不要任何其他文字）：
{
  "formula": "=SUMIFS(金额列, 部门列, A2, 月份列, B2)",
  "description": "公式用途的一句话说明",
  "params": [
    {"name": "参数名或单元格引用", "meaning": "含义与取值规则", "example": "示例取值"}
  ],
  "logic": ["计算步骤1", "计算步骤2"],
  "scenarios": "适用场景、数据前置条件与版本要求",
  "alternatives": [
    {"formula": "替代公式表达式", "scenario": "适用场景/优缺点"}
  ],
  "pitfalls": "常见使用陷阱与避坑建议"
}

要求：
- formula 必须是可直接粘贴到单元格的完整表达式
- params 覆盖公式中每一个引用范围/条件（逐项解释）
- 用通俗语言解释，让非技术用户也能理解"""
            prompt = data.data.get("prompt", "")
            result = call_llm(system_prompt, prompt)
        elif data.operation == "outliers":
            raw = data.data.get("raw", "") or data.data.get("content", "")
            result = json.dumps(_detect_outliers(raw), ensure_ascii=False)
        else:
            result = json.dumps({"status": "created", "data": data.data})

        conn.execute(
            "INSERT INTO excel_operations (id, operation, title, data, result) VALUES (?,?,?,?,?)",
            (op_id, data.operation, data.title, json.dumps(data.data), result),
        )
        conn.commit()
        return {"ok": True, "id": op_id, "result": result}
    except Exception as e:
        raise HTTPException(500, "操作失败，请稍后重试") from e
    finally:
        conn.close()


@router.get("/api/excel/history")
async def list_excel_history(current_user: dict = require_auth()):
    conn = get_db()
    try:
        where, args = _user_scope_clause(conn, current_user)
        items = []
        for row in conn.execute(
            f"SELECT * FROM excel_operations WHERE 1=1{where} ORDER BY created_at DESC LIMIT 50", args
        ).fetchall():
            e = dict(row)
            e["data"] = json.loads(e.get("data", "{}"))
            items.append(e)
        return items
    finally:
        conn.close()


# ── 内容创作商业化升级：异步任务处理器注册（进度/自动重试/并发控制）──
register_handler("copywriting_generate", _copywriting_handler, user_limit=2, max_attempts=1)
register_handler("translation_translate", _translation_handler, user_limit=2, max_attempts=1)
register_handler("ppt_generate", _ppt_handler, user_limit=2, max_attempts=1)


def _setup_test_environment(project_dir: str, cfg: dict) -> bool:
    """准备测试环境：安装依赖、检查测试文件。"""
    try:
        # 检查测试文件
        test_files = []
        for root, dirs, files in os.walk(project_dir):
            dirs[:] = [d for d in dirs if d not in ('.venv', 'venv', '__pycache__', 'node_modules')]
            for f in files:
                if f.endswith('_test.py') or f.startswith('test_') or f.endswith('.test.js'):
                    test_files.append(os.path.join(root, f))
        
        if not test_files:
            logger.warning("未找到测试文件")
            return False
        
        # 安装依赖
        req_file = os.path.join(project_dir, 'requirements.txt')
        if os.path.exists(req_file):
            subprocess.run(['pip', 'install', '-r', req_file], 
                         cwd=project_dir, capture_output=True, timeout=120)
        
        return True
    except Exception as e:
        logger.error(f"准备测试环境失败: {e}")
        return False


def _run_core_tests(project_dir: str, test_cmd: list, timeout: int = 300) -> tuple:
    """运行核心测试。"""
    try:
        result = subprocess.run(test_cmd, 
                               cwd=project_dir,
                               capture_output=True,
                               text=True,
                               timeout=timeout)
        return result.returncode, result.stdout, result.stderr
    except subprocess.TimeoutExpired:
        return 124, "", "测试超时"
    except Exception as e:
        return 1, "", str(e)


def _validate_results(returncode: int, stdout: str, stderr: str) -> dict:
    """验证测试结果。"""
    results = {
        'passed': 0,
        'failed': 0,
        'error': 0,
        'skipped': 0,
        'summary': ''
    }
    
    # 解析 pytest 输出
    import re
    passed_match = re.search(r'(\d+) passed', stdout)
    failed_match = re.search(r'(\d+) failed', stdout)
    error_match = re.search(r'(\d+) error', stdout)
    skipped_match = re.search(r'(\d+) skipped', stdout)
    
    if passed_match:
        results['passed'] = int(passed_match.group(1))
    if failed_match:
        results['failed'] = int(failed_match.group(1))
    if error_match:
        results['error'] = int(error_match.group(1))
    if skipped_match:
        results['skipped'] = int(skipped_match.group(1))
    
    # 生成摘要
    if results['failed'] == 0 and results['error'] == 0:
        results['summary'] = '✅ 所有测试通过'
    else:
        results['summary'] = f'❌ {results["failed"]} 个失败, {results["error"]} 个错误'
    
    return results
